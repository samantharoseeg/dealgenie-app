"""
DealGenie Pro - Commercial Real Estate Analysis Platform
Main Streamlit Application
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
from datetime import datetime, timedelta
import json
import re
from typing import Dict, Any, Optional, List, Tuple
import base64
import io
from PIL import Image
import io
import os
from pathlib import Path
from ocr_parser import ComprehensiveDataParser
from llm_enhancement import render_api_settings, render_summary_with_llm_option, calculate_metrics_for_llm
from cre_extraction_engine import CREExtractionEngine, ASSET_CLASSES
from benchmarks import (
    BENCHMARKS as BENCHMARK_DATA,
    METRICS_CATALOG,
    get_benchmark_range,
    get_status,
    get_metric_info,
    get_all_metrics_for_asset_class
)
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
import matplotlib.pyplot as plt
import seaborn as sns
import xlsxwriter

# Page Configuration
st.set_page_config(
    page_title="DealGenie Pro | CRE Analysis Platform",
    page_icon="🏢",
    layout="wide",
    initial_sidebar_state="collapsed",
    menu_items={
        'Get Help': 'https://github.com/yourusername/dealgenie-app',
        'Report a bug': "https://github.com/yourusername/dealgenie-app/issues",
        'About': "DealGenie Pro v1.0 - Institutional-Grade CRE Analysis"
    }
)

# ============================================================================
# TEMPLATE MANAGEMENT SYSTEM
# ============================================================================

# Template directory path
TEMPLATES_DIR = Path("data/templates")

def ensure_templates_dir():
    """Ensure templates directory exists"""
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

def list_templates() -> List[str]:
    """Return list of available template names"""
    ensure_templates_dir()
    template_files = list(TEMPLATES_DIR.glob("*.json"))
    return [f.stem for f in template_files]

def save_template(template_name: str, settings_dict: Dict) -> bool:
    """
    Save template to JSON file

    Args:
        template_name: Name of the template
        settings_dict: Dictionary containing benchmark_overrides and custom_dd_items

    Returns:
        True if successful, False otherwise
    """
    try:
        ensure_templates_dir()
        template_path = TEMPLATES_DIR / f"{template_name}.json"

        # Prepare data to save
        template_data = {
            "template_name": template_name,
            "created_date": datetime.now().isoformat(),
            "benchmark_overrides": settings_dict.get("benchmark_overrides", {}),
            "custom_dd_items": settings_dict.get("custom_dd_items", {}),
            "profile_name": settings_dict.get("profile_name", "")
        }

        # Write to file
        with open(template_path, 'w') as f:
            json.dump(template_data, f, indent=2)

        return True
    except Exception as e:
        st.error(f"Error saving template: {str(e)}")
        return False

def load_template(template_name: str) -> bool:
    """
    Load template from JSON file and update session state

    Args:
        template_name: Name of the template to load

    Returns:
        True if successful, False otherwise
    """
    try:
        template_path = TEMPLATES_DIR / f"{template_name}.json"

        if not template_path.exists():
            st.error(f"Template '{template_name}' not found")
            return False

        # Read template file
        with open(template_path, 'r') as f:
            template_data = json.load(f)

        # Update session state
        if "benchmark_overrides" in template_data:
            st.session_state.benchmark_overrides = template_data["benchmark_overrides"]

        if "custom_dd_items" in template_data:
            st.session_state.custom_dd_items = template_data["custom_dd_items"]

        if "profile_name" in template_data:
            st.session_state.profile_name = template_data["profile_name"]

        # Update template tracking flags
        st.session_state.template_loaded = True
        st.session_state.active_template_name = template_name
        st.session_state.active_template_date = template_data.get("created_date", "Unknown")

        return True
    except Exception as e:
        st.error(f"Error loading template: {str(e)}")
        return False

def delete_template(template_name: str) -> bool:
    """
    Delete a template file

    Args:
        template_name: Name of the template to delete

    Returns:
        True if successful, False otherwise
    """
    try:
        template_path = TEMPLATES_DIR / f"{template_name}.json"
        if template_path.exists():
            template_path.unlink()
            return True
        return False
    except Exception as e:
        st.error(f"Error deleting template: {str(e)}")
        return False

def validate_template_structure(template_data: Dict) -> Tuple[bool, str]:
    """
    Validate imported template JSON structure

    Args:
        template_data: Dictionary loaded from JSON file

    Returns:
        Tuple of (is_valid, error_message)
    """
    try:
        # Check required keys
        required_keys = ["template_name", "created_date"]
        for key in required_keys:
            if key not in template_data:
                return False, f"Missing required key: '{key}'"

        # Validate benchmark_overrides structure (if present)
        if "benchmark_overrides" in template_data:
            if not isinstance(template_data["benchmark_overrides"], dict):
                return False, "benchmark_overrides must be a dictionary"
            # Check nested structure
            for asset_class, subclasses in template_data["benchmark_overrides"].items():
                if not isinstance(subclasses, dict):
                    return False, f"Invalid structure in benchmark_overrides for '{asset_class}'"
                for subclass, metrics in subclasses.items():
                    if not isinstance(metrics, dict):
                        return False, f"Invalid structure in benchmark_overrides for '{asset_class}.{subclass}'"

        # Validate custom_dd_items structure (if present)
        if "custom_dd_items" in template_data:
            if not isinstance(template_data["custom_dd_items"], dict):
                return False, "custom_dd_items must be a dictionary"
            # Check that values are lists
            for category, items in template_data["custom_dd_items"].items():
                if not isinstance(items, list):
                    return False, f"custom_dd_items['{category}'] must be a list"
                # Check that list contains strings
                for item in items:
                    if not isinstance(item, str):
                        return False, f"Items in custom_dd_items['{category}'] must be strings"

        # Validate profile_name (if present)
        if "profile_name" in template_data:
            if not isinstance(template_data["profile_name"], str):
                return False, "profile_name must be a string"

        return True, ""

    except Exception as e:
        return False, f"Validation error: {str(e)}"

def import_template_from_json(json_data: Dict) -> bool:
    """
    Import template from JSON data and update session state

    Args:
        json_data: Dictionary containing template data

    Returns:
        True if successful, False otherwise
    """
    try:
        # Validate structure
        is_valid, error_msg = validate_template_structure(json_data)
        if not is_valid:
            st.error(f"❌ Invalid template format: {error_msg}")
            return False

        # Update session state with validated data
        if "benchmark_overrides" in json_data:
            st.session_state.benchmark_overrides = json_data["benchmark_overrides"]

        if "custom_dd_items" in json_data:
            st.session_state.custom_dd_items = json_data["custom_dd_items"]

        if "profile_name" in json_data:
            st.session_state.profile_name = json_data["profile_name"]

        # Store template metadata for display
        st.session_state.imported_template_name = json_data.get("template_name", "Imported Template")
        st.session_state.imported_template_date = json_data.get("created_date", "Unknown")

        return True

    except Exception as e:
        st.error(f"❌ Error importing template: {str(e)}")
        return False

def export_template_to_json() -> Dict:
    """
    Export current settings as JSON-serializable dictionary

    Returns:
        Dictionary containing all template data
    """
    # Get current analysis data to determine asset class defaults
    asset_class_defaults = {}
    if hasattr(st.session_state, 'analysis_data') and st.session_state.analysis_data:
        asset_class_defaults = {
            "asset_class": st.session_state.analysis_data.get("asset_class", ""),
            "subclass": st.session_state.analysis_data.get("subclass", "")
        }

    template_data = {
        "template_name": st.session_state.get("profile_name", "Exported Template"),
        "created_date": datetime.now().isoformat(),
        "export_date": datetime.now().isoformat(),
        "benchmark_overrides": st.session_state.get("benchmark_overrides", {}),
        "custom_dd_items": st.session_state.get("custom_dd_items", {}),
        "profile_name": st.session_state.get("profile_name", ""),
        "asset_class_defaults": asset_class_defaults,
        "metadata": {
            "version": "1.0",
            "exported_from": "DealGenie Pro"
        }
    }

    return template_data

def get_default_template() -> Optional[str]:
    """
    Get the name of the default template

    Returns:
        Template name or None if no default set
    """
    try:
        default_file = Path("data/default_template.txt")
        if default_file.exists():
            return default_file.read_text().strip()
        return None
    except Exception as e:
        return None

def set_default_template(template_name: str) -> bool:
    """
    Set a template as the default

    Args:
        template_name: Name of template to set as default

    Returns:
        True if successful
    """
    try:
        default_file = Path("data/default_template.txt")
        default_file.parent.mkdir(parents=True, exist_ok=True)
        default_file.write_text(template_name)
        return True
    except Exception as e:
        st.error(f"Error setting default template: {str(e)}")
        return False

def clear_default_template() -> bool:
    """
    Clear the default template setting

    Returns:
        True if successful
    """
    try:
        default_file = Path("data/default_template.txt")
        if default_file.exists():
            default_file.unlink()
        return True
    except Exception as e:
        st.error(f"Error clearing default template: {str(e)}")
        return False

def duplicate_template(template_name: str, new_name: str) -> bool:
    """
    Duplicate an existing template

    Args:
        template_name: Name of template to duplicate
        new_name: Name for the new copy

    Returns:
        True if successful
    """
    try:
        source_path = TEMPLATES_DIR / f"{template_name}.json"
        if not source_path.exists():
            st.error(f"Template '{template_name}' not found")
            return False

        # Read source template
        with open(source_path, 'r') as f:
            template_data = json.load(f)

        # Update name and date
        template_data["template_name"] = new_name
        template_data["created_date"] = datetime.now().isoformat()
        if "original_template" not in template_data:
            template_data["original_template"] = template_name

        # Save as new template
        return save_template(new_name, template_data)

    except Exception as e:
        st.error(f"Error duplicating template: {str(e)}")
        return False

def get_template_metadata(template_name: str) -> Optional[Dict]:
    """
    Get metadata for a template

    Args:
        template_name: Name of template

    Returns:
        Dictionary with metadata or None
    """
    try:
        template_path = TEMPLATES_DIR / f"{template_name}.json"
        if not template_path.exists():
            return None

        with open(template_path, 'r') as f:
            template_data = json.load(f)

        # Count customizations
        num_benchmarks = 0
        if "benchmark_overrides" in template_data:
            for asset_class, subclasses in template_data["benchmark_overrides"].items():
                for subclass, metrics in subclasses.items():
                    num_benchmarks += len(metrics)

        num_dd_items = sum(len(items) for items in template_data.get("custom_dd_items", {}).values())

        return {
            "name": template_name,
            "created_date": template_data.get("created_date", "Unknown"),
            "num_benchmarks": num_benchmarks,
            "num_dd_items": num_dd_items,
            "profile_name": template_data.get("profile_name", ""),
            "description": template_data.get("description", ""),
            "base_asset_class": template_data.get("asset_class_defaults", {}).get("asset_class", "")
        }

    except Exception as e:
        return None

def compare_templates(template1_name: str, template2_name: str) -> Dict:
    """
    Compare two templates and show differences

    Args:
        template1_name: Name of first template
        template2_name: Name of second template

    Returns:
        Dictionary with comparison data
    """
    try:
        # Load both templates
        path1 = TEMPLATES_DIR / f"{template1_name}.json"
        path2 = TEMPLATES_DIR / f"{template2_name}.json"

        if not path1.exists() or not path2.exists():
            return {"error": "One or both templates not found"}

        with open(path1, 'r') as f:
            data1 = json.load(f)
        with open(path2, 'r') as f:
            data2 = json.load(f)

        # Compare benchmark overrides
        benchmarks1 = data1.get("benchmark_overrides", {})
        benchmarks2 = data2.get("benchmark_overrides", {})

        # Count differences
        all_asset_classes = set(benchmarks1.keys()) | set(benchmarks2.keys())
        benchmark_diffs = []

        for asset_class in all_asset_classes:
            if asset_class not in benchmarks1:
                benchmark_diffs.append(f"➕ {asset_class} (only in {template2_name})")
            elif asset_class not in benchmarks2:
                benchmark_diffs.append(f"➖ {asset_class} (only in {template1_name})")

        # Compare DD items
        dd1 = data1.get("custom_dd_items", {})
        dd2 = data2.get("custom_dd_items", {})

        all_categories = set(dd1.keys()) | set(dd2.keys())
        dd_diffs = []

        for category in all_categories:
            items1 = set(dd1.get(category, []))
            items2 = set(dd2.get(category, []))

            only_in_1 = items1 - items2
            only_in_2 = items2 - items1

            if only_in_1:
                dd_diffs.append(f"➖ {category}: {len(only_in_1)} items only in {template1_name}")
            if only_in_2:
                dd_diffs.append(f"➕ {category}: {len(only_in_2)} items only in {template2_name}")

        return {
            "template1": template1_name,
            "template2": template2_name,
            "benchmark_diffs": benchmark_diffs if benchmark_diffs else ["No differences"],
            "dd_diffs": dd_diffs if dd_diffs else ["No differences"]
        }

    except Exception as e:
        return {"error": f"Comparison error: {str(e)}"}

# ============================================================================
# CUSTOM STYLING
# ============================================================================

def create_metric_help_text(metric: str) -> str:
    """Create help text for a metric using the catalog"""
    info = get_metric_info(metric)
    return f"{info.get('description', '')}\n\n{info.get('why_it_matters', '')}"

def inject_custom_css():
    """Apply custom CSS styling for professional look"""
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

    .stApp {
        font-family: 'Inter', sans-serif;
        background: linear-gradient(180deg, #f8f9fa 0%, #ffffff 100%);
    }

    .main-header {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 16px;
        margin-bottom: 2rem;
        box-shadow: 0 10px 30px rgba(0,0,0,0.1);
        color: white;
    }

    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 12px;
        box-shadow: 0 2px 8px rgba(0,0,0,0.08);
        transition: all 0.3s ease;
        height: 100%;
    }

    .metric-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    }

    .status-badge {
        padding: 0.375rem 0.875rem;
        border-radius: 20px;
        font-size: 0.813rem;
        font-weight: 600;
        display: inline-block;
        text-transform: uppercase;
    }

    .status-good { background: #10b981; color: white; }
    .status-warning { background: #f59e0b; color: white; }
    .status-critical { background: #ef4444; color: white; }

    @media (max-width: 768px) {
        .main-header { padding: 1rem; }
        .metric-card { margin-bottom: 1rem; }
    }
    </style>
    """, unsafe_allow_html=True)

# ============================================================================
# FINANCIAL CALCULATIONS
# ============================================================================

def calculate_mortgage_constant(annual_rate: float, amort_years: int, io_period: int = 0) -> float:
    """Calculate mortgage constant with IO period support"""
    if amort_years == 0:
        return annual_rate

    monthly_rate = annual_rate / 12
    n_payments = amort_years * 12

    if monthly_rate == 0:
        monthly_payment = 1 / n_payments
    else:
        monthly_payment = (monthly_rate * (1 + monthly_rate)**n_payments) / ((1 + monthly_rate)**n_payments - 1)

    return monthly_payment * 12

def get_metric_info(metric_name: str) -> Dict:
    """Get metric information from METRICS_CATALOG"""
    try:
        from benchmarks import METRICS_CATALOG
        return METRICS_CATALOG.get(metric_name, {
            "unit": "",
            "description": f"Metric: {metric_name}",
            "why_it_matters": "Important for investment analysis"
        })
    except ImportError:
        return {
            "unit": "",
            "description": f"Metric: {metric_name}",
            "why_it_matters": "Important for investment analysis"
        }

def get_all_metrics_for_asset_class(asset_class: str, subclass: str) -> Dict[str, List[str]]:
    """
    Get filtered metrics relevant to specific asset class and subclass
    Returns categorized metrics that should be displayed
    """
    # Define asset-specific metric mappings
    ASSET_METRIC_MAP = {
        "multifamily": {
            "all": {  # Common to all multifamily subclasses
                "leverage": ["ltv", "dscr", "debt_yield", "interest_rate"],
                "valuation": ["cap_rate", "exit_cap", "price_per_unit"],
                "operating": ["occupancy", "expense_ratio", "noi_growth", "replacement_reserves"],
                "revenue": ["avg_rent", "market_rent", "rent_growth", "concessions"],
                "physical": ["units", "avg_unit_size", "parking_ratio"]
            },
            "garden_lowrise": {},
            "midrise": {},
            "highrise": {"additional": ["concierge_cost", "amenity_ratio"]},
            "student": {"additional": ["beds_per_unit", "summer_occupancy"]},
            "senior": {"additional": ["care_level", "medicare_mix"]}
        },
        "office": {
            "all": {
                "leverage": ["ltv", "dscr", "debt_yield", "interest_rate"],
                "valuation": ["cap_rate", "exit_cap", "price_psf"],
                "operating": ["occupancy", "expense_ratio", "noi_growth"],
                "leasing": ["walt", "tenant_improvement", "leasing_commission", "renewal_probability"],
                "physical": ["gla_sf", "floor_plate", "parking_ratio"]
            },
            "cbd_a_trophy": {},
            "cbd_b_commodity": {},
            "suburban": {},
            "medical": {"additional": ["hospital_affiliation", "medicaid_mix"]},
            "life_sciences": {"additional": ["lab_space_pct", "power_density"]}
        },
        "industrial": {
            "all": {
                "leverage": ["ltv", "dscr", "debt_yield", "interest_rate"],
                "valuation": ["cap_rate", "exit_cap", "price_psf"],
                "operating": ["occupancy", "expense_ratio", "noi_growth"],
                "physical": ["clear_height", "dock_doors", "drive_in_doors", "power_density"],
                "location": ["distance_to_port", "distance_to_airport", "highway_access"]
            },
            "bulk_warehouse": {"additional": ["cross_dock_capable", "trailer_parking"]},
            "light_industrial": {"additional": ["office_finish_pct"]},
            "flex": {"additional": ["office_percentage", "loading_ratio"]},
            "cold_storage": {"additional": ["temperature_zones", "refrigeration_redundancy"]},
            "data_center": {"additional": ["power_capacity_mw", "cooling_redundancy", "fiber_connectivity"]}
        },
        "retail": {
            "all": {
                "leverage": ["ltv", "dscr", "debt_yield", "interest_rate"],
                "valuation": ["cap_rate", "exit_cap", "price_psf"],
                "operating": ["occupancy", "expense_ratio", "noi_growth"],
                "sales": ["sales_psf", "rent_to_sales"],
                "tenant": ["anchor_tenant", "anchor_remaining_term"]
            },
            "power_center": {},
            "neighborhood_center": {"additional": ["grocery_anchor"]},
            "lifestyle_center": {"additional": ["restaurant_pct", "entertainment_pct"]},
            "single_tenant_nnn": {"additional": ["lease_term_remaining", "rent_escalations"]}
        },
        "hospitality": {
            "all": {
                "leverage": ["ltv", "dscr", "debt_yield", "interest_rate"],
                "valuation": ["cap_rate", "exit_cap", "price_per_key"],
                "operating": ["occupancy", "adr", "revpar", "gop_margin"],
                "physical": ["keys", "meeting_space_sf"],
                "brand": ["brand_flag", "franchise_fee", "pip_required"]
            },
            "limited_service": {},
            "full_service": {"additional": ["banquet_capacity", "room_service"]},
            "luxury": {"additional": ["michelin_restaurants", "butler_service"]},
            "resort": {"additional": ["golf_courses", "seasonality"]}
        }
    }

    # Get metrics for the asset class
    if asset_class not in ASSET_METRIC_MAP:
        return {
            "💰 Leverage & Debt": ["ltv", "dscr", "debt_yield", "interest_rate"],
            "📊 Valuation": ["cap_rate", "exit_cap"],
            "🏢 Operating": ["occupancy", "expense_ratio", "noi_growth"]
        }

    asset_metrics = ASSET_METRIC_MAP[asset_class]
    result = {}

    # Add common metrics
    if "all" in asset_metrics:
        for category, metrics in asset_metrics["all"].items():
            display_category = {
                "leverage": "💰 Leverage & Debt",
                "valuation": "📊 Valuation & Returns",
                "operating": "🏢 Operating Metrics",
                "revenue": "💵 Revenue Metrics",
                "physical": "🏗️ Physical Characteristics",
                "leasing": "📝 Leasing Metrics",
                "tenant": "👥 Tenant Metrics",
                "location": "📍 Location Factors",
                "sales": "💳 Sales Performance",
                "brand": "🏨 Brand & Franchise"
            }.get(category, f"📋 {category.title()}")
            result[display_category] = metrics

    # Add subclass-specific metrics
    if subclass in asset_metrics and "additional" in asset_metrics[subclass]:
        result["🎯 Asset-Specific"] = asset_metrics[subclass]["additional"]

    # Filter by available benchmarks if they exist
    if asset_class in BENCHMARK_DATA and subclass in BENCHMARK_DATA[asset_class]:
        available_benchmarks = set(BENCHMARK_DATA[asset_class][subclass].keys())
        filtered_result = {}
        for category, metrics in result.items():
            filtered_metrics = [m for m in metrics if m in available_benchmarks]
            if filtered_metrics:
                filtered_result[category] = filtered_metrics
        return filtered_result

    return result

def calculate_dscr(noi: float, loan_amount: float, rate: float, amort_years: int) -> float:
    """Calculate Debt Service Coverage Ratio"""
    if loan_amount == 0 or rate == 0:
        return 0

    mortgage_constant = calculate_mortgage_constant(rate, amort_years)
    annual_debt_service = loan_amount * mortgage_constant

    return noi / annual_debt_service if annual_debt_service > 0 else 0

def calculate_irr(cash_flows: List[float]) -> float:
    """Calculate Internal Rate of Return"""
    try:
        return np.irr(cash_flows) * 100
    except:
        return 0

# ============================================================================
# OCR PARSER
# ============================================================================

# Legacy parser for backward compatibility
class FinancialDataParser:
    """Legacy parser - redirects to ComprehensiveDataParser"""

    def parse(self, text: str) -> Dict[str, Any]:
        """Use comprehensive parser but return simplified format for compatibility"""
        parser = ComprehensiveDataParser()
        result = parser.parse(text)

        # Extract key fields for backward compatibility
        fields = result.get('extracted_fields', {})
        simplified = {
            'confidence': result.get('overall_confidence', 0.8),
            'purchase_price': fields.get('purchase_price'),
            'noi': fields.get('noi'),
            'cap_rate': fields.get('cap_rate'),
            'loan_amount': fields.get('loan_amount'),
            'interest_rate': fields.get('interest_rate'),
            'units': fields.get('unit_count'),
            'sf': fields.get('building_sf')
        }

        # Remove None values
        return {k: v for k, v in simplified.items() if v is not None}

# ============================================================================
# BENCHMARKS WITH SOURCES
# ============================================================================

BENCHMARKS = {
    "Office": {
        "cap_rate": {"min": 5.5, "preferred": 6.5, "max": 7.5, "source": "CBRE Q4 2024"},
        "dscr": {"min": 1.25, "preferred": 1.40, "max": 1.60, "source": "MBA Survey Q4 2024"},
        "ltv": {"min": 50, "preferred": 65, "max": 75, "source": "CMBS Market Report 2024"}
    },
    "Multifamily": {
        "cap_rate": {"min": 4.5, "preferred": 5.5, "max": 6.5, "source": "RCA Analytics Q4 2024"},
        "dscr": {"min": 1.20, "preferred": 1.35, "max": 1.50, "source": "Freddie Mac 2024"},
        "ltv": {"min": 55, "preferred": 70, "max": 80, "source": "Fannie Mae Guidelines 2024"}
    },
    "Industrial": {
        "cap_rate": {"min": 5.0, "preferred": 6.0, "max": 7.0, "source": "JLL Research Q4 2024"},
        "dscr": {"min": 1.25, "preferred": 1.40, "max": 1.55, "source": "Life Co Survey 2024"},
        "ltv": {"min": 55, "preferred": 65, "max": 70, "source": "CMBS Market Report 2024"}
    },
    "Retail": {
        "cap_rate": {"min": 6.0, "preferred": 7.0, "max": 8.0, "source": "Cushman & Wakefield 2024"},
        "dscr": {"min": 1.35, "preferred": 1.50, "max": 1.75, "source": "Regional Banks Survey"},
        "ltv": {"min": 50, "preferred": 60, "max": 65, "source": "Insurance Co Guidelines"}
    },
    "Hotel": {
        "cap_rate": {"min": 7.0, "preferred": 8.5, "max": 10.0, "source": "STR/HVS Report 2024"},
        "dscr": {"min": 1.30, "preferred": 1.45, "max": 1.60, "source": "Hospitality Lenders 2024"},
        "ltv": {"min": 50, "preferred": 60, "max": 65, "source": "CMBS Hotel Loans 2024"}
    }
}

def evaluate_against_benchmarks(asset_class: str, metrics: Dict) -> List[Dict]:
    """Evaluate metrics against industry benchmarks"""
    if asset_class not in BENCHMARKS:
        return []

    benchmarks = BENCHMARKS[asset_class]
    evaluations = []

    for metric, value in metrics.items():
        if metric in benchmarks:
            bench = benchmarks[metric]
            status = "good"
            if value < bench["min"]:
                status = "critical"
            elif value < bench["preferred"]:
                status = "warning"

            evaluations.append({
                "metric": metric.upper(),
                "value": value,
                "status": status,
                "benchmark": f"{bench['preferred']} ({bench['source']})"
            })

    return evaluations

# ============================================================================
# MAIN APPLICATION
# ============================================================================

def _display_fields(extracted_data: Dict, field_names: List[str]):
    """Helper function to display extracted fields in a clean format"""
    for field in field_names:
        if field in extracted_data:
            value = extracted_data[field]
            # Format the field name
            display_name = field.replace('_', ' ').title()

            # Format the value
            if isinstance(value, float):
                if field.endswith('_pct') or field == 'cap_rate':
                    st.write(f"**{display_name}**: {value:.2f}%")
                elif 'price' in field or 'cost' in field or 'amount' in field:
                    st.write(f"**{display_name}**: ${value:,.0f}")
                else:
                    st.write(f"**{display_name}**: {value:,.2f}")
            elif isinstance(value, list):
                st.write(f"**{display_name}**:")
                for item in value[:5]:  # Show first 5 items
                    if isinstance(item, dict):
                        st.caption(f"  • {item}")
                    else:
                        st.caption(f"  • {item}")
            else:
                st.write(f"**{display_name}**: {value}")

def render_asset_specific_fields(asset_class: str, subclass: str = None) -> Dict:
    """
    Render asset-specific input fields based on property type
    Returns dictionary of field values
    """
    specific_data = {}

    st.markdown("### 🏗️ Asset-Specific Details")

    if asset_class == "Office":
        # Get relevant metrics for office and filter display
        asset_class_lower = "office"
        relevant_metrics = get_all_metrics_for_asset_class(asset_class_lower, subclass or "suburban")
        all_relevant_metrics = set()
        for category_metrics in relevant_metrics.values():
            all_relevant_metrics.update(category_metrics)

        col1, col2, col3 = st.columns(3)

        with col1:
            # Always show GLA for office
            if 'gla_sf' in all_relevant_metrics:
                gla_sf = st.number_input(
                    "GLA (SF)",
                    min_value=0,
                    value=st.session_state.get('office_gla', 100000),
                    step=1000,
                    key='office_gla',
                    help="Gross Leasable Area in square feet"
                )
            else:
                gla_sf = 100000

            # Show WALT for office
            if 'walt' in all_relevant_metrics:
                walt_years = st.number_input(
                    "WALT (years)",
                    min_value=0.0,
                    max_value=20.0,
                    value=st.session_state.get('office_walt', 5.5),
                    step=0.5,
                    key='office_walt',
                    help="Weighted Average Lease Term"
                )
            else:
                walt_years = 5.5

            # Show tenant count
            tenant_count = st.number_input(
                "Tenant Count",
                min_value=1,
                value=st.session_state.get('office_tenant_count', 10),
                step=1,
                key='office_tenant_count'
            )

        with col2:
            # Show TI only if in relevant metrics
            if 'tenant_improvement' in all_relevant_metrics:
                ti_new_psf = st.number_input(
                    "TI New ($/SF)",
                    min_value=0.0,
                    value=st.session_state.get('office_ti_new', 75.0),
                    step=5.0,
                    key='office_ti_new',
                    help="Tenant Improvement allowance for new leases"
                )
                ti_renewal_psf = st.number_input(
                    "TI Renewal ($/SF)",
                    min_value=0.0,
                    value=st.session_state.get('office_ti_renewal', 25.0),
                    step=5.0,
                    key='office_ti_renewal',
                    help="Tenant Improvement allowance for renewals"
                )
            else:
                ti_new_psf = 0.0
                ti_renewal_psf = 0.0

            top5_tenants_pct = st.number_input(
                "Top 5 Tenants (%)",
                min_value=0.0,
                max_value=100.0,
                value=st.session_state.get('office_top5_pct', 60.0),
                step=5.0,
                key='office_top5_pct',
                help="Percentage of rent from top 5 tenants"
            )

        with col3:
            # Show LC only if in relevant metrics
            if 'leasing_commission' in all_relevant_metrics:
                lc_new_pct = st.number_input(
                    "LC New (%)",
                    min_value=0.0,
                    max_value=10.0,
                    value=st.session_state.get('office_lc_new', 5.5),
                    step=0.5,
                    key='office_lc_new',
                    help="Leasing Commission for new leases"
                )
                lc_renewal_pct = st.number_input(
                    "LC Renewal (%)",
                    min_value=0.0,
                    max_value=10.0,
                    value=st.session_state.get('office_lc_renewal', 2.5),
                    step=0.5,
                    key='office_lc_renewal',
                    help="Leasing Commission for renewals"
                )
            else:
                lc_new_pct = 0.0
                lc_renewal_pct = 0.0

        specific_data = {
            'gla_sf': gla_sf,
            'walt_years': walt_years,
            'ti_new_psf': ti_new_psf,
            'ti_renewal_psf': ti_renewal_psf,
            'lc_new_pct': lc_new_pct / 100,
            'lc_renewal_pct': lc_renewal_pct / 100,
            'tenant_count': tenant_count,
            'top5_tenants_pct': top5_tenants_pct / 100
        }

    elif asset_class == "Multifamily":
        col1, col2, col3 = st.columns(3)

        with col1:
            units = st.number_input(
                "Units",
                min_value=1,
                value=st.session_state.get('mf_units', 200),
                step=1,
                key='mf_units',
                help="Total number of apartment units"
            )
            avg_rent = st.number_input(
                "Avg Rent ($)",
                min_value=0,
                value=st.session_state.get('mf_avg_rent', 1500),
                step=50,
                key='mf_avg_rent',
                help="Average monthly rent per unit"
            )

        with col2:
            market_rent = st.number_input(
                "Market Rent ($)",
                min_value=0,
                value=st.session_state.get('mf_market_rent', 1650),
                step=50,
                key='mf_market_rent',
                help="Market monthly rent per unit"
            )
            occupancy_pct = st.number_input(
                "Occupancy (%)",
                min_value=0.0,
                max_value=100.0,
                value=st.session_state.get('mf_occupancy', 94.0),
                step=1.0,
                key='mf_occupancy'
            )

        with col3:
            expense_ratio = st.number_input(
                "Expense Ratio (%)",
                min_value=0.0,
                max_value=100.0,
                value=st.session_state.get('mf_expense_ratio', 40.0),
                step=1.0,
                key='mf_expense_ratio',
                help="Operating expenses as % of revenue"
            )
            concessions_months = st.number_input(
                "Concessions (months)",
                min_value=0.0,
                max_value=6.0,
                value=st.session_state.get('mf_concessions', 1.0),
                step=0.5,
                key='mf_concessions',
                help="Free rent concessions in months"
            )

        specific_data = {
            'units': units,
            'avg_rent': avg_rent,
            'market_rent': market_rent,
            'occupancy_pct': occupancy_pct / 100,
            'expense_ratio': expense_ratio / 100,
            'concessions_months': concessions_months
        }

    elif asset_class == "Retail":
        col1, col2, col3 = st.columns(3)

        with col1:
            gla_sf = st.number_input(
                "GLA (SF)",
                min_value=0,
                value=st.session_state.get('retail_gla', 75000),
                step=1000,
                key='retail_gla',
                help="Gross Leasable Area"
            )
            anchor_tenant = st.text_input(
                "Anchor Tenant",
                value=st.session_state.get('retail_anchor', 'Kroger'),
                key='retail_anchor',
                help="Primary anchor tenant name"
            )

        with col2:
            anchor_term_years = st.number_input(
                "Anchor Term Remaining (years)",
                min_value=0.0,
                max_value=30.0,
                value=st.session_state.get('retail_anchor_term', 12.0),
                step=0.5,
                key='retail_anchor_term'
            )
            sales_psf = st.number_input(
                "Sales PSF ($)",
                min_value=0.0,
                value=st.session_state.get('retail_sales_psf', 450.0),
                step=25.0,
                key='retail_sales_psf',
                help="Average tenant sales per square foot"
            )

        with col3:
            co_tenancy = st.checkbox(
                "Co-Tenancy Clause",
                value=st.session_state.get('retail_co_tenancy', True),
                key='retail_co_tenancy',
                help="Does the lease have co-tenancy provisions?"
            )

        specific_data = {
            'gla_sf': gla_sf,
            'anchor_tenant': anchor_tenant,
            'anchor_term_years': anchor_term_years,
            'co_tenancy_clause': co_tenancy,
            'sales_psf': sales_psf
        }

    elif asset_class == "Industrial":
        # Get relevant metrics for industrial and filter display
        asset_class_lower = "industrial"
        relevant_metrics = get_all_metrics_for_asset_class(asset_class_lower, subclass or "bulk_warehouse")
        all_relevant_metrics = set()
        for category_metrics in relevant_metrics.values():
            all_relevant_metrics.update(category_metrics)

        col1, col2, col3 = st.columns(3)

        with col1:
            building_sf = st.number_input(
                "Building SF",
                min_value=0,
                value=st.session_state.get('ind_building_sf', 150000),
                step=5000,
                key='ind_building_sf'
            )

            # Always show clear height for industrial
            if 'clear_height' in all_relevant_metrics:
                clear_height = st.number_input(
                    "Clear Height (ft)",
                min_value=0,
                max_value=60,
                value=st.session_state.get('ind_clear_height', 32),
                step=1,
                key='ind_clear_height',
                help="Clear ceiling height in feet"
            )

            else:
                clear_height = 32  # Default if not relevant

        with col2:
            # Only show dock doors for warehouse types
            if 'dock_doors' in all_relevant_metrics:
                dock_doors = st.number_input(
                    "Dock Doors (count)",
                    min_value=0,
                    value=st.session_state.get('ind_dock_doors', 20),
                    step=1,
                    key='ind_dock_doors'
                )
            else:
                dock_doors = 0

            # Show office finish for flex/light industrial
            if 'office_finish_pct' in all_relevant_metrics or 'office_percentage' in all_relevant_metrics:
                office_finish_pct = st.number_input(
                    "Office Finish (%)",
                    min_value=0.0,
                    max_value=100.0,
                    value=st.session_state.get('ind_office_finish', 8.0),
                    step=1.0,
                    key='ind_office_finish',
                    help="Percentage of building that is office space"
                )
            else:
                office_finish_pct = 0

        with col3:
            # Only show cold storage for cold storage subclass
            if subclass == "cold_storage":
                temperature_zones = st.number_input(
                    "Temperature Zones",
                    min_value=1,
                    max_value=5,
                    value=st.session_state.get('ind_temp_zones', 2),
                    step=1,
                    key='ind_temp_zones',
                    help="Number of different temperature zones"
                )
                cold_storage = True
            else:
                temperature_zones = 0
                cold_storage = False

            # Show power density for data centers
            if subclass == "data_center" and 'power_capacity_mw' in all_relevant_metrics:
                power_capacity = st.number_input(
                    "Power Capacity (MW)",
                    min_value=0.0,
                    value=st.session_state.get('ind_power_mw', 5.0),
                    step=0.5,
                    key='ind_power_mw',
                    help="Total power capacity in megawatts"
                )
            else:
                power_capacity = 0

        specific_data = {
            'building_sf': building_sf,
            'clear_height_ft': clear_height if 'clear_height' in all_relevant_metrics else None,
            'dock_doors': dock_doors if 'dock_doors' in all_relevant_metrics else None,
            'office_finish_pct': office_finish_pct / 100 if office_finish_pct > 0 else None,
            'cold_storage': cold_storage,
            'temperature_zones': temperature_zones if temperature_zones > 0 else None,
            'power_capacity_mw': power_capacity if power_capacity > 0 else None
        }
        # Remove None values
        specific_data = {k: v for k, v in specific_data.items() if v is not None}

    elif asset_class == "Hotel" or asset_class == "Hospitality":
        # Get relevant metrics for hospitality and filter display
        asset_class_lower = "hospitality"
        relevant_metrics = get_all_metrics_for_asset_class(asset_class_lower, subclass or "limited_service")
        all_relevant_metrics = set()
        for category_metrics in relevant_metrics.values():
            all_relevant_metrics.update(category_metrics)

        col1, col2, col3 = st.columns(3)

        with col1:
            keys = st.number_input(
                "Keys",
                min_value=1,
                value=st.session_state.get('hotel_keys', 120),
                step=1,
                key='hotel_keys',
                help="Number of hotel rooms/keys"
            )
            adr = st.number_input(
                "ADR ($)",
                min_value=0.0,
                value=st.session_state.get('hotel_adr', 150.0),
                step=10.0,
                key='hotel_adr',
                help="Average Daily Rate"
            )
            brand_flag = st.text_input(
                "Brand/Flag",
                value=st.session_state.get('hotel_brand', 'Marriott'),
                key='hotel_brand',
                help="Hotel brand or flag"
            )

        with col2:
            occupancy_pct = st.number_input(
                "Occupancy (%)",
                min_value=0.0,
                max_value=100.0,
                value=st.session_state.get('hotel_occupancy', 72.0),
                step=1.0,
                key='hotel_occupancy'
            )
            revpar = st.number_input(
                "RevPAR ($)",
                min_value=0.0,
                value=st.session_state.get('hotel_revpar', 108.0),
                step=5.0,
                key='hotel_revpar',
                help="Revenue Per Available Room"
            )

        with col3:
            gop_margin_pct = st.number_input(
                "GOP Margin (%)",
                min_value=0.0,
                max_value=100.0,
                value=st.session_state.get('hotel_gop_margin', 38.0),
                step=1.0,
                key='hotel_gop_margin',
                help="Gross Operating Profit margin"
            )
            pip_cost_per_key = st.number_input(
                "PIP Cost per Key ($)",
                min_value=0,
                value=st.session_state.get('hotel_pip_cost', 15000),
                step=1000,
                key='hotel_pip_cost',
                help="Property Improvement Plan cost per key"
            )

        specific_data = {
            'keys': keys,
            'adr': adr,
            'occupancy_pct': occupancy_pct / 100,
            'revpar': revpar,
            'gop_margin_pct': gop_margin_pct / 100,
            'brand_flag': brand_flag,
            'pip_cost_per_key': pip_cost_per_key
        }

    return specific_data

def render_header():
    """Render application header"""
    st.markdown("""
    <div class="main-header">
        <h1 style="margin: 0; font-size: 2.5rem;">🏢 DealGenie Pro</h1>
        <p style="margin: 0.5rem 0 0 0; opacity: 0.9;">Institutional-Grade Commercial Real Estate Analysis Platform</p>
    </div>
    """, unsafe_allow_html=True)

def render_input_section():
    """Render deal input section"""
    st.header("📊 Deal Analysis")

    # Input method selection
    input_method = st.radio(
        "Select Input Method",
        ["📝 Manual Entry", "📷 Photo/OCR", "📁 Upload File"],
        horizontal=True
    )

    parsed_data = {}

    if input_method == "📝 Manual Entry":
        # Define asset class to subclass mapping
        ASSET_CLASSES = {
            "Multifamily": ["garden_lowrise", "midrise", "highrise", "student_housing", "senior_IL", "senior_AL", "senior_MemoryCare"],
            "Office": ["cbd_A_trophy", "cbd_BC", "suburban", "medical_office", "flex_creative"],
            "Industrial": ["bulk_distribution", "manufacturing", "cold_storage", "last_mile"],
            "Retail": ["neighborhood_center", "community_center", "power_center", "lifestyle_center", "outlet_center"],
            "Hotel": ["full_service_upscale", "select_service", "extended_stay", "resort", "limited_service"]
        }

        col1, col2, col3 = st.columns(3)

        with col1:
            asset_class = st.selectbox(
                "Asset Class",
                ["Office", "Multifamily", "Industrial", "Retail", "Hotel"]
            )

            # Add subclass dropdown
            subclass_options = ASSET_CLASSES.get(asset_class, [])
            if subclass_options:
                subclass = st.selectbox(
                    "Property Type",
                    subclass_options,
                    format_func=lambda x: x.replace("_", " ").title()
                )
                st.session_state['subclass'] = subclass

            purchase_price = st.number_input(
                "Purchase Price ($)",
                min_value=0,
                value=18500000,
                step=100000
            )

        with col2:
            noi = st.number_input(
                "Year 1 NOI ($)",
                min_value=0,
                value=1110000,
                step=10000
            )
            loan_amount = st.number_input(
                "Loan Amount ($)",
                min_value=0,
                value=13000000,
                step=100000
            )

        with col3:
            interest_rate = st.slider(
                "Interest Rate (%)",
                min_value=3.0,
                max_value=10.0,
                value=6.5,
                step=0.25
            ) / 100
            amort_years = st.number_input(
                "Amortization (years)",
                min_value=0,
                max_value=40,
                value=30
            )

        # Add asset-specific fields after basic deal fields
        st.markdown("---")
        specific_fields = render_asset_specific_fields(asset_class)

        # Combine basic data with asset-specific data
        parsed_data = {
            "asset_class": asset_class,
            "purchase_price": purchase_price,
            "noi": noi,
            "loan_amount": loan_amount,
            "interest_rate": interest_rate,
            "amort_years": amort_years,
            **specific_fields  # Merge asset-specific fields
        }

    elif input_method == "📷 Photo/OCR":
        # Asset Class Selection for Enhanced Extraction
        st.markdown("### 🏢 Asset Classification")

        col1, col2 = st.columns(2)
        with col1:
            asset_class = st.selectbox(
                "Select Asset Class",
                options=list(ASSET_CLASSES.keys()),
                help="Choose the property type for enhanced extraction and benchmarks"
            )

        with col2:
            if asset_class in ASSET_CLASSES:
                subclass = st.selectbox(
                    "Select Subclass",
                    options=ASSET_CLASSES[asset_class],
                    help="Choose specific subclass for targeted analysis"
                )
            else:
                subclass = None

        # Store selections in session state
        if 'asset_class' not in st.session_state:
            st.session_state.asset_class = asset_class
        if 'subclass' not in st.session_state:
            st.session_state.subclass = subclass

        st.session_state.asset_class = asset_class
        st.session_state.subclass = subclass

        st.markdown("---")

        uploaded_file = st.file_uploader(
            "Upload deal summary (Image, PDF, or PowerPoint)",
            type=['png', 'jpg', 'jpeg', 'pdf', 'pptx', 'ppt']
        )

        if uploaded_file:
            file_type = uploaded_file.name.split('.')[-1].lower()
            ocr_text = ""

            # Process based on file type
            if file_type in ['png', 'jpg', 'jpeg']:
                st.info("📸 Processing image with OCR...")
                try:
                    import pytesseract
                    image = Image.open(uploaded_file)
                    ocr_text = pytesseract.image_to_string(image)

                    # Check if no text was extracted
                    if not ocr_text.strip():
                        st.warning("⚠️ No text found in image. The image may be unclear or contain no readable text. Using demo data.")
                        ocr_text = """
                        INVESTMENT SUMMARY
                        Property: Northgate Business Center
                        Address: 1234 Market Street, Dallas, TX 75201
                        Year Built: 1985 | Renovated: 2020
                        Building Size: 125,000 SF
                        Site: 5.2 acres
                        Parking: 350 spaces (2.8/1,000 SF)
                        Occupancy: 92%
                        WALT: 4.2 years
                        Number of Tenants: 12
                        Anchor Tenant: Wells Fargo (25,000 SF)

                        FINANCIAL HIGHLIGHTS
                        Purchase Price: $18.5MM
                        Price/SF: $148
                        NOI: $1,110,000
                        Cap Rate: 6.0%
                        T-12 EGI: $1,850,000
                        Operating Expenses: $740,000
                        Real Estate Taxes: $285,000
                        Insurance: $48,000
                        Management Fee: 3.5%

                        DEBT TERMS
                        Loan Amount: $13 million
                        LTV: 70%
                        Interest Rate: SOFR + 250 bps (6.25% all-in)
                        Amortization: 30 years
                        IO Period: 3 years
                        Term: 10 years
                        DSCR Requirement: 1.25x minimum
                        Origination Fee: 1.0%
                        Extension: 2x12mo at 0.25% fee
                        Rate Cap: 7.5% strike

                        EXIT STRATEGY
                        Hold Period: 5 years
                        Exit Cap Rate: 6.75%
                        Disposition Fee: 1.5%
                        """

                except Exception as e:
                    st.warning("OCR library not available. Using demo data.")
                    # Fallback to demo text if OCR fails
                    ocr_text = """
                    Purchase Price: $18.5MM
                    NOI: $1,110,000
                    Cap Rate: 6.0%
                    Loan Amount: $13 million
                    Interest Rate: 6.25%
                    """

            elif file_type == 'pdf':
                st.info("📄 Processing PDF document...")
                try:
                    import pdfplumber
                    with pdfplumber.open(uploaded_file) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                ocr_text += page_text + "\n"

                    # Check if no text was extracted
                    if not ocr_text.strip():
                        st.warning("⚠️ No text found in PDF. The file may be image-based or protected. Using demo data.")
                        ocr_text = """
                        INVESTMENT SUMMARY
                        Property: Northgate Business Center
                        Address: 1234 Market Street, Dallas, TX 75201
                        Year Built: 1985 | Renovated: 2020
                        Building Size: 125,000 SF
                        Site: 5.2 acres
                        Parking: 350 spaces (2.8/1,000 SF)
                        Occupancy: 92%
                        WALT: 4.2 years
                        Number of Tenants: 12
                        Anchor Tenant: Wells Fargo (25,000 SF)

                        FINANCIAL HIGHLIGHTS
                        Purchase Price: $18.5MM
                        Price/SF: $148
                        NOI: $1,110,000
                        Cap Rate: 6.0%
                        T-12 EGI: $1,850,000
                        Operating Expenses: $740,000
                        Real Estate Taxes: $285,000
                        Insurance: $48,000
                        Management Fee: 3.5%

                        DEBT TERMS
                        Loan Amount: $13 million
                        LTV: 70%
                        Interest Rate: SOFR + 250 bps (6.25% all-in)
                        Amortization: 30 years
                        IO Period: 3 years
                        Term: 10 years
                        DSCR Requirement: 1.25x minimum
                        Origination Fee: 1.0%
                        Extension: 2x12mo at 0.25% fee
                        Rate Cap: 7.5% strike

                        EXIT STRATEGY
                        Hold Period: 5 years
                        Exit Cap Rate: 6.75%
                        Disposition Fee: 1.5%
                        """

                except Exception as e:
                    st.warning("PDF processing library not available. Using demo data.")
                    ocr_text = """
                    Purchase Price: $18.5MM
                    NOI: $1,110,000
                    Cap Rate: 6.0%
                    Loan Amount: $13 million
                    Interest Rate: 6.25%
                    """

            elif file_type in ['pptx', 'ppt']:
                st.info("📊 Processing PowerPoint presentation...")
                try:
                    from pptx import Presentation
                    prs = Presentation(uploaded_file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                ocr_text += shape.text + "\n"

                    # Check if no text was extracted
                    if not ocr_text.strip():
                        st.warning("⚠️ No text found in PowerPoint. The slides may contain only images or shapes. Using demo data.")
                        ocr_text = """
                        INVESTMENT SUMMARY
                        Property: Northgate Business Center
                        Address: 1234 Market Street, Dallas, TX 75201
                        Year Built: 1985 | Renovated: 2020
                        Building Size: 125,000 SF
                        Site: 5.2 acres
                        Parking: 350 spaces (2.8/1,000 SF)
                        Occupancy: 92%
                        WALT: 4.2 years
                        Number of Tenants: 12
                        Anchor Tenant: Wells Fargo (25,000 SF)

                        FINANCIAL HIGHLIGHTS
                        Purchase Price: $18.5MM
                        Price/SF: $148
                        NOI: $1,110,000
                        Cap Rate: 6.0%
                        T-12 EGI: $1,850,000
                        Operating Expenses: $740,000
                        Real Estate Taxes: $285,000
                        Insurance: $48,000
                        Management Fee: 3.5%

                        DEBT TERMS
                        Loan Amount: $13 million
                        LTV: 70%
                        Interest Rate: SOFR + 250 bps (6.25% all-in)
                        Amortization: 30 years
                        IO Period: 3 years
                        Term: 10 years
                        DSCR Requirement: 1.25x minimum
                        Origination Fee: 1.0%
                        Extension: 2x12mo at 0.25% fee
                        Rate Cap: 7.5% strike

                        EXIT STRATEGY
                        Hold Period: 5 years
                        Exit Cap Rate: 6.75%
                        Disposition Fee: 1.5%
                        """

                except Exception as e:
                    st.warning("PowerPoint processing library not available. Using demo data.")
                    ocr_text = """
                    Purchase Price: $18.5MM
                    NOI: $1,110,000
                    Cap Rate: 6.0%
                    Loan Amount: $13 million
                    Interest Rate: 6.25%
                    """

            else:
                st.error("Unsupported file type")
                return parsed_data

            if ocr_text.strip():
                # Use enhanced CRE extraction engine if asset class selected
                if hasattr(st.session_state, 'asset_class') and hasattr(st.session_state, 'subclass'):
                    try:
                        st.info(f"🚀 Using Enhanced Extraction for {st.session_state.asset_class.title()} - {st.session_state.subclass.replace('_', ' ').title()}")

                        # Import and call extract_and_analyze function
                        from cre_extraction_engine import extract_and_analyze, load_benchmarks

                        # Get benchmark library
                        benchmark_library = load_benchmarks()

                        # Create OCR blocks structure if we can detect tables in text
                        ocr_blocks = []
                        # Simple table detection - look for lines with consistent delimiters
                        lines = ocr_text.split('\n')
                        for i, line in enumerate(lines):
                            # Check if line looks like table data (has multiple columns)
                            if any(delimiter in line for delimiter in ['|', '\t', '  ']) and len(line.split()) > 2:
                                ocr_blocks.append({
                                    'text': line,
                                    'type': 'table_row',
                                    'line': i,
                                    'bbox': {'x': 0, 'y': i * 20}  # Simple positioning
                                })
                            # Check for headers
                            elif any(header in line.upper() for header in ['METRIC', 'VALUE', 'FINANCIAL', 'DEBT', 'TERMS']):
                                ocr_blocks.append({
                                    'text': line,
                                    'type': 'header',
                                    'line': i,
                                    'bbox': {'x': 0, 'y': i * 20}
                                })

                        # Get benchmark overrides for this asset class/subclass
                        benchmark_overrides = None
                        if 'benchmark_overrides' in st.session_state:
                            asset_class = st.session_state.get('asset_class')
                            subclass = st.session_state.get('subclass')
                            if (asset_class in st.session_state.benchmark_overrides and
                                subclass in st.session_state.benchmark_overrides[asset_class]):
                                benchmark_overrides = st.session_state.benchmark_overrides[asset_class][subclass]

                        # Call the extraction and analysis function with OCR blocks and overrides
                        extracted_data = extract_and_analyze(
                            asset_class=st.session_state.get('asset_class'),
                            subclass=st.session_state.get('subclass'),
                            raw_text=ocr_text,
                            benchmark_library=benchmark_library,
                            ocr_blocks=ocr_blocks if ocr_blocks else None,
                            benchmark_overrides=benchmark_overrides
                        )

                        # Store results in session state
                        st.session_state['extracted_data'] = extracted_data
                        st.session_state['field_confidence'] = extracted_data.get('confidence', {})
                        st.session_state['parsed_data'] = extracted_data.get('ingested', {})

                        # For backward compatibility, also store as cre_result
                        st.session_state.cre_result = extracted_data

                        # Display enhanced results
                        col1, col2 = st.columns([2, 1])

                        with col1:
                            # Show extracted text in expander
                            with st.expander("📝 Extracted Text", expanded=False):
                                st.text(ocr_text[:2000] + "..." if len(ocr_text) > 2000 else ocr_text)

                        with col2:
                            # Show extraction summary with new metrics
                            completeness = extracted_data.get('completeness', {})
                            required_fields = completeness.get('required_fields', 0)
                            total_required = completeness.get('total_required', 1)
                            completeness_pct = (required_fields / total_required) * 100 if total_required > 0 else 0

                            st.metric("Completeness", f"{completeness_pct:.0f}%")
                            st.caption(f"Required: {required_fields}/{total_required}")
                            st.caption(f"Total fields: {len(extracted_data.get('ingested', {}))}")

                            # Show glossary terms if found
                            glossary_refs = extracted_data.get('glossary_refs', [])
                            if glossary_refs:
                                st.caption(f"📚 Terms: {', '.join(glossary_refs[:5])}")

                            # Show errors if any
                            if 'error' in extracted_data:
                                st.error(extracted_data['error'])

                    except Exception as e:
                        st.error(f"Enhanced extraction failed: {str(e)}")
                        st.info("Falling back to basic parser...")

                        # Fallback to comprehensive parser
                        comprehensive_parser = ComprehensiveDataParser()
                        comprehensive_result = comprehensive_parser.parse(ocr_text)

                        # Show extraction results
                        col1, col2 = st.columns([2, 1])

                        with col1:
                            # Show extracted text in expander
                            with st.expander("📝 Extracted Text", expanded=False):
                                st.text(ocr_text[:2000] + "..." if len(ocr_text) > 2000 else ocr_text)

                        with col2:
                            # Show extraction summary
                            st.metric(
                                "Extraction Confidence",
                                f"{comprehensive_result['overall_confidence']*100:.0f}%"
                            )
                            st.caption(f"Fields found: {len(comprehensive_result['extracted_fields'])}")

                else:
                    # Use basic comprehensive parser if no asset class selected
                    comprehensive_parser = ComprehensiveDataParser()
                    comprehensive_result = comprehensive_parser.parse(ocr_text)

                    # Show extraction results
                    col1, col2 = st.columns([2, 1])

                    with col1:
                        # Show extracted text in expander
                        with st.expander("📝 Extracted Text", expanded=False):
                            st.text(ocr_text[:2000] + "..." if len(ocr_text) > 2000 else ocr_text)

                    with col2:
                        # Show extraction summary
                        st.metric(
                            "Extraction Confidence",
                            f"{comprehensive_result['overall_confidence']*100:.0f}%"
                        )
                        st.caption(f"Fields found: {len(comprehensive_result['extracted_fields'])}")

                # Display extracted data - Enhanced CRE results or fallback
                if hasattr(st.session_state, 'extracted_data') and st.session_state.extracted_data:
                    # Use extracted_data from extract_and_analyze
                    cre_result = st.session_state.extracted_data
                elif hasattr(st.session_state, 'cre_result') and st.session_state.cre_result:
                    # Display enhanced CRE extraction results
                    cre_result = st.session_state.cre_result

                    st.subheader("🚀 Enhanced CRE Analysis")

                    # Enhanced display with tabs for new structure
                    tabs = st.tabs(["📋 Ingested", "📊 Derived", "⚖️ Benchmarks", "⚠️ Risks", "📈 Sensitivities", "❓ Unknown"])

                    with tabs[0]:  # Ingested Fields
                        st.markdown("**Extracted Fields**")
                        if cre_result.get('ingested'):
                            for field, value in cre_result['ingested'].items():
                                confidence = cre_result.get('confidence', {}).get(field, 'Medium')

                                # Format value for display
                                if isinstance(value, float):
                                    if field.endswith('_pct') or field.endswith('_rate'):
                                        formatted_value = f"{value:.2%}"
                                    elif value > 1000:
                                        formatted_value = f"${value:,.0f}"
                                    else:
                                        formatted_value = f"{value:.3f}"
                                else:
                                    formatted_value = str(value)

                                st.write(f"**{field.replace('_', ' ').title()}:** {formatted_value} ({confidence} confidence)")

                    with tabs[1]:  # Derived Metrics
                        st.markdown("**Computed Metrics**")
                        if cre_result.get('derived'):
                            for metric, value in cre_result['derived'].items():
                                if not metric.endswith('_calc'):
                                    # Format derived values
                                    if isinstance(value, float):
                                        if metric in ['cap_rate', 'yield_on_cost', 'debt_yield', 'dscr']:
                                            formatted_value = f"{value:.3f}"
                                        elif metric in ['exit_value', 'net_sale_proceeds', 'refi_proceeds']:
                                            formatted_value = f"${value:,.0f}"
                                        else:
                                            formatted_value = f"{value:.3f}"
                                    else:
                                        formatted_value = str(value)

                                    # Show calculation if available
                                    calc = cre_result['derived'].get(f"{metric}_calc", "")
                                    st.write(f"**{metric.replace('_', ' ').title()}:** {formatted_value}")
                                    if calc:
                                        st.caption(f"Calculation: {calc}")

                    with tabs[2]:  # Benchmark Comparisons
                        st.markdown("**Benchmark Analysis**")
                        if cre_result.get('bench_compare'):
                            for field, comparison in cre_result['bench_compare'].items():
                                status = comparison.get('status', 'Unknown')
                                benchmark_info = comparison.get('benchmark', 'N/A')
                                source = comparison.get('source', 'Industry Research')

                                if status == 'OK':
                                    st.success(f"✅ **{field.replace('_', ' ').title()}:** {status}")
                                elif status in ['Above Target', 'Offside High']:
                                    st.info(f"📈 **{field.replace('_', ' ').title()}:** {status}")
                                elif status in ['Below Target', 'Offside Low']:
                                    st.warning(f"📉 **{field.replace('_', ' ').title()}:** {status}")
                                elif status == 'Poor':
                                    st.error(f"🔴 **{field.replace('_', ' ').title()}:** {status}")
                                else:
                                    st.write(f"**{field.replace('_', ' ').title()}:** {status}")

                                # Show benchmark source
                                st.caption(f"Benchmark: {benchmark_info}")

                    with tabs[3]:  # Risk Analysis
                        st.markdown("**Risk Assessment**")
                        if cre_result.get('risks_ranked'):
                            for risk in cre_result['risks_ranked']:
                                severity = risk.get('severity', 'Medium')
                                if severity == 'High':
                                    st.error(f"🔴 **{risk['metric'].replace('_', ' ').title()}:** {risk['issue']}")
                                elif severity == 'Medium':
                                    st.warning(f"🟡 **{risk['metric'].replace('_', ' ').title()}:** {risk['issue']}")
                                else:
                                    st.info(f"🔵 **{risk['metric'].replace('_', ' ').title()}:** {risk['issue']}")

                                # Show mitigations
                                if risk.get('mitigations'):
                                    st.caption("Mitigations:")
                                    for mitigation in risk['mitigations']:
                                        st.caption(f"• {mitigation}")

                    with tabs[4]:  # Sensitivity Analysis
                        st.markdown("**Sensitivity Analysis**")
                        if cre_result.get('sensitivities'):
                            for metric, scenarios in cre_result['sensitivities'].items():
                                st.write(f"**{metric.replace('_', ' ').title()} Sensitivity:**")
                                for scenario, values in scenarios.items():
                                    st.caption(f"• {scenario}: " + ", ".join([f"{k}={v:.2f}" if isinstance(v, float) else f"{k}={v}" for k, v in values.items()]))

                    with tabs[5]:  # Unknown Fields
                        st.markdown("**Missing Data**")
                        if cre_result.get('unknown'):
                            for item in cre_result['unknown']:
                                st.write(f"• {item}")

                        # Display notes if available
                        if cre_result.get('notes'):
                            st.markdown("**Processing Notes**")
                            for note in cre_result['notes'][:10]:
                                st.caption(f"📝 {note}")

                elif 'comprehensive_result' in locals() and comprehensive_result['extracted_fields']:
                    # Fallback to basic comprehensive parser display
                    st.subheader("📊 Extracted Data")

                    tabs = st.tabs(["Deal Info", "Financials", "Debt Terms", "Operations", "Development"])

                    with tabs[0]:  # Deal Info
                        deal_fields = ['property_name', 'street_address', 'city', 'state', 'zip_code',
                                      'year_built', 'year_renovated', 'building_sf', 'unit_count',
                                      'parking_spaces', 'occupancy_pct', 'walt_years', 'anchor_tenant']
                        _display_fields(comprehensive_result['extracted_fields'], deal_fields)

                    with tabs[1]:  # Financials
                        financial_fields = ['purchase_price', 'noi', 'cap_rate', 'gross_income',
                                          'operating_expenses', 'exit_cap_rate', 'hold_period_years',
                                          'closing_costs', 'disposition_fee_pct']
                        _display_fields(comprehensive_result['extracted_fields'], financial_fields)

                    with tabs[2]:  # Debt Terms
                        debt_fields = ['loan_amount', 'interest_rate', 'amort_years', 'io_period_years',
                                      'loan_term_years', 'ltv_pct', 'min_dscr', 'origination_fee_pct',
                                      'rate_cap_strike', 'extension_count', 'extension_term_months']
                        _display_fields(comprehensive_result['extracted_fields'], debt_fields)

                    with tabs[3]:  # Operations
                        ops_fields = ['real_estate_taxes', 'insurance_cost', 'management_fee_pct',
                                     'replacement_reserves', 'vacancy_rate', 'market_rent',
                                     'ti_allowance_new', 'leasing_commission_pct', 'free_rent_months']
                        _display_fields(comprehensive_result['extracted_fields'], ops_fields)

                    with tabs[4]:  # Development
                        dev_fields = ['land_cost', 'hard_costs', 'soft_costs', 'developer_fee',
                                     'contingency_pct', 'preleasing_pct', 'expected_delivery',
                                     'construction_contract_type', 'general_contractor']
                        _display_fields(comprehensive_result['extracted_fields'], dev_fields)

                    # Show extraction notes
                    if comprehensive_result['extraction_notes']:
                        with st.expander("📋 Extraction Notes", expanded=False):
                            for note in comprehensive_result['extraction_notes'][:10]:  # Show first 10
                                st.caption(f"• {note}")

                    # Show missing critical fields
                    if comprehensive_result['missing_critical']:
                        st.warning(f"⚠️ Missing critical fields: {', '.join(comprehensive_result['missing_critical'])}")

                # Convert to legacy format for compatibility
                if hasattr(st.session_state, 'cre_result') and st.session_state.cre_result:
                    # Use enhanced CRE results
                    cre_result = st.session_state.cre_result
                    parsed_data = {
                        'confidence': 0.9,  # Higher confidence for enhanced extraction
                        'purchase_price': cre_result.get('ingested', {}).get('purchase_price'),
                        'noi': cre_result.get('ingested', {}).get('noi_now'),
                        'cap_rate': cre_result.get('derived', {}).get('cap_rate') or cre_result.get('ingested', {}).get('entry_cap'),
                        'loan_amount': cre_result.get('ingested', {}).get('loan_amount'),
                        'interest_rate': cre_result.get('ingested', {}).get('rate'),
                        'asset_class': st.session_state.asset_class.title(),
                        'subclass': st.session_state.subclass
                    }

                    # Add more fields from ingested data
                    ingested = cre_result.get('ingested', {})
                    for key in ['ltv', 'dscr', 'occupancy_pct', 'expense_ratio']:
                        if key in ingested:
                            parsed_data[key] = ingested[key]

                    # Add derived metrics
                    derived = cre_result.get('derived', {})
                    for key in ['equity_multiple', 'irr', 'yield_on_cost']:
                        if key in derived:
                            parsed_data[key] = derived[key]

                    completeness = cre_result.get('completeness', {})
                    required_fields = completeness.get('required_fields', 0)
                    total_required = completeness.get('total_required', 1)
                    st.success(f"🚀 Enhanced extraction completed: {required_fields}/{total_required} required fields ({len(cre_result.get('ingested', {}))+ len(cre_result.get('derived', {}))} total)")

                else:
                    # Fallback to basic parser
                    parser = FinancialDataParser()
                    parsed_data = parser.parse(ocr_text)
                    parsed_data["asset_class"] = "Office"  # Default

                    if 'comprehensive_result' in locals():
                        st.success(f"✅ Extracted {len(comprehensive_result['extracted_fields'])} fields with {comprehensive_result['overall_confidence']*100:.0f}% confidence")

    return parsed_data

def generate_principal_summary(data: Dict) -> str:
    """
    Generate principal-style investment summary from extracted data

    Args:
        data: Either legacy data dict or extracted_data from CRE extraction engine

    Returns:
        Formatted investment summary with benchmark citations
    """

    # Check if we have enhanced extracted_data or legacy data
    if isinstance(data, dict) and 'ingested' in data and 'bench_compare' in data:
        # Use enhanced extraction data
        return generate_enhanced_principal_summary(data)
    else:
        # Fall back to legacy summary generation for backward compatibility
        return generate_legacy_principal_summary(data)


def generate_enhanced_principal_summary(extracted_data: Dict) -> str:
    """
    Generate enhanced principal summary using CRE extraction engine output

    Args:
        extracted_data: Output from extract_and_analyze function

    Returns:
        Structured investment summary with three sections: Strengths, Concerns, Verdict
    """

    sections = []

    # SECTION 1: STRENGTHS (from benchmark comparisons)
    strengths = []
    bench_compare = extracted_data.get('bench_compare', {})

    for metric, comparison in bench_compare.items():
        status = comparison.get('status', '')
        if status in ['OK', 'Above Target']:
            value = comparison.get('value')
            benchmark = comparison.get('benchmark', '')
            source = comparison.get('source', 'Industry Benchmark')

            # Format the metric name nicely
            metric_name = metric.replace('_', ' ').upper()

            # Format value based on metric type
            if metric in ['cap_rate', 'dscr', 'debt_yield']:
                formatted_value = f"{value:.2f}" if isinstance(value, (int, float)) else str(value)
            elif metric in ['ltv', 'occupancy_pct', 'expense_ratio']:
                formatted_value = f"{value:.1f}%" if isinstance(value, (int, float)) else str(value)
            else:
                formatted_value = str(value)

            # Build strength statement with source
            if benchmark and source != 'Industry Benchmark':
                strengths.append(f"{metric_name} at {formatted_value} {benchmark}")
            else:
                strengths.append(f"{metric_name} at {formatted_value} within target range ({source})")

    # Add derived metric strengths
    derived = extracted_data.get('derived', {})
    if 'equity_multiple' in derived and derived['equity_multiple'] > 1.8:
        strengths.append(f"Equity multiple of {derived['equity_multiple']:.2f}x exceeds 1.8x target")
    if 'irr' in derived and derived['irr'] > 15:
        strengths.append(f"Projected IRR of {derived['irr']:.1f}% surpasses 15% hurdle")

    if strengths:
        sections.append("**STRENGTHS:**\n" + "\n".join([f"• {s}" for s in strengths[:5]]))  # Limit to 5

    # SECTION 2: CONCERNS (from risk ranking)
    concerns = []
    risks_ranked = extracted_data.get('risks_ranked', [])

    for risk in risks_ranked:
        if risk.get('severity') == 'High':
            issue = risk.get('issue', risk.get('metric', 'Unknown issue'))
            mitigations = risk.get('mitigations', [])

            concern_text = f"Issue: {issue}."
            if mitigations:
                # Take first mitigation and add dollar amount if present
                mitigation = mitigations[0]
                if '$' in mitigation or 'cost' in mitigation.lower():
                    concern_text += f" To fix: {mitigation}"
                else:
                    concern_text += f" To fix: {mitigation} (est. $50-100k)"

            concerns.append(concern_text)

    # Add missing critical fields as concerns
    unknown = extracted_data.get('unknown', [])
    for item in unknown[:3]:  # Limit to 3 missing items
        if 'critical' in item.lower() or 'required' in item.lower():
            concerns.append(f"Missing: {item}")

    if not concerns:
        # Add standard concerns if none found
        if 'exit_cap' not in extracted_data.get('ingested', {}):
            concerns.append("Issue: Exit cap assumption not specified. To fix: Conduct market study ($5-10k)")
        if derived.get('dscr', 0) < 1.3:
            concerns.append("Issue: Thin debt coverage. To fix: Negotiate rate reduction or increase equity")

    if concerns:
        sections.append("**CONCERNS:**\n" + "\n".join([f"• {c}" for c in concerns[:4]]))  # Limit to 4

    # SECTION 3: VERDICT (based on risk analysis)
    high_risks = len([r for r in risks_ranked if r.get('severity') == 'High'])
    medium_risks = len([r for r in risks_ranked if r.get('severity') == 'Medium'])
    completeness = extracted_data.get('completeness', {})
    completeness_score = completeness.get('percent', 0) if isinstance(completeness, dict) else 0

    # Determine verdict
    if high_risks >= 3:
        verdict = "**REJECT** - Multiple high-severity risks exceed acceptable threshold"
    elif high_risks >= 2 and completeness_score < 70:
        verdict = "**REVIEW REQUIRED** - Significant risks with incomplete data require further diligence"
    elif high_risks == 1 and medium_risks <= 2:
        verdict = "**PROCEED WITH CONDITIONS** - Address primary risk before closing"
    elif completeness_score >= 80 and high_risks == 0:
        verdict = "**PASS** - Strong fundamentals with manageable risk profile justify proceeding"
    else:
        verdict = "**REVIEW** - Moderate risk profile requires IC discussion"

    sections.append(f"**VERDICT:** {verdict}")

    # Combine sections
    summary = "\n\n".join(sections)

    # Apply LLM enhancement if configured
    summary = apply_llm_enhancement_to_summary(summary, extracted_data)

    return summary


def generate_legacy_principal_summary(data: Dict) -> str:
    """
    Legacy principal summary generation for backward compatibility
    Keeps the original logic when extraction engine data is not available
    """

    # Calculate key metrics
    dscr = calculate_dscr(
        data.get('noi', 0),
        data.get('loan_amount', 0),
        data.get('interest_rate', 0.065),
        data.get('amort_years', 30)
    )

    cap_rate = (data.get('noi', 0) / data.get('purchase_price', 1)) * 100 if data.get('purchase_price', 0) > 0 else 0
    ltv = (data.get('loan_amount', 0) / data.get('purchase_price', 1)) * 100 if data.get('purchase_price', 0) > 0 else 0

    # Simple equity multiple calculation
    equity = data.get('purchase_price', 0) - data.get('loan_amount', 0)
    exit_cap = data.get('exit_cap_rate', cap_rate + 0.5)
    hold_period = data.get('hold_period', 5)

    if exit_cap > 0 and hold_period > 0 and equity > 0:
        future_noi = data.get('noi', 0) * (1.03 ** hold_period)
        exit_value = future_noi / (exit_cap / 100)
        net_proceeds = exit_value - data.get('loan_amount', 0) * 0.9
        equity_multiple = net_proceeds / equity if equity > 0 else 0
        irr = ((equity_multiple ** (1/hold_period)) - 1) * 100 if equity_multiple > 0 else 0
    else:
        equity_multiple = 0
        irr = 0

    # Build structured summary
    sections = []

    # STRENGTHS
    strengths = []
    if dscr > 1.3:
        strengths.append(f"DSCR at {dscr:.2f}x exceeds minimum 1.25x (MBA Survey Q4 2024)")
    if irr > 15:
        strengths.append(f"IRR of {irr:.1f}% surpasses 15% hurdle rate")
    if ltv < 65:
        strengths.append(f"Conservative {ltv:.0f}% LTV provides refinancing flexibility (CMBS Market Report 2024)")
    if cap_rate > 6:
        strengths.append(f"Entry cap of {cap_rate:.2f}% above market average (CBRE Q4 2024)")

    if strengths:
        sections.append("**STRENGTHS:**\n" + "\n".join([f"• {s}" for s in strengths]))

    # CONCERNS
    concerns = []
    if equity_multiple < 1.6:
        concerns.append(f"Issue: {equity_multiple:.2f}x equity multiple below 1.6x target. To fix: Extend hold period or reduce leverage")
    if exit_cap > cap_rate + 0.5:
        concerns.append(f"Issue: Exit cap expansion to {exit_cap:.1f}% pressures returns. To fix: Stress test at 50bps higher")
    if dscr < 1.25:
        concerns.append(f"Issue: Thin {dscr:.2f}x coverage. To fix: Negotiate interest-only period ($0 cost)")

    if concerns:
        sections.append("**CONCERNS:**\n" + "\n".join([f"• {c}" for c in concerns]))

    # VERDICT
    if equity_multiple >= 1.8 and irr >= 15:
        verdict = "**PASS** - Strong risk-adjusted returns justify proceeding"
    elif equity_multiple >= 1.5 and dscr >= 1.25:
        verdict = "**PROCEED WITH CONDITIONS** - Workable with modest leverage reduction"
    elif dscr < 1.2 or equity_multiple < 1.3:
        verdict = "**REJECT** - Insufficient margin of safety"
    else:
        verdict = "**REVIEW** - Marginal returns require IC discussion"

    sections.append(f"**VERDICT:** {verdict}")

    # Combine sections
    summary = "\n\n".join(sections)

    # Apply LLM enhancement if configured
    summary = apply_llm_enhancement_to_summary(summary, data)

    return summary


def apply_llm_enhancement_to_summary(summary: str, data: Dict) -> str:
    """
    Apply LLM polish to the structured summary if API is configured

    Args:
        summary: The structured summary text
        data: Original data for context

    Returns:
        Polished summary or original if no API configured
    """

    # Check if LLM API is configured
    if not hasattr(st.session_state, 'api_provider') or not st.session_state.api_provider:
        return summary

    if not hasattr(st.session_state, 'api_key') or not st.session_state.api_key:
        return summary

    provider = st.session_state.api_provider
    api_key = st.session_state.api_key

    # Prepare the prompt
    prompt = f"""Polish this CRE investment summary while preserving ALL numbers, sources, and technical terms.
Keep it concise and principal-focused. Maintain the STRENGTHS, CONCERNS, VERDICT structure.
Do not add any new information or analysis.

Original Summary:
{summary}

Output only the polished summary."""

    try:
        if "Claude" in provider:
            import anthropic
            client = anthropic.Anthropic(api_key=api_key)
            response = client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=500,
                temperature=0.3,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text.strip()

        elif "OpenAI" in provider:
            from openai import OpenAI
            client = OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a senior CRE investment principal. Write concise, decisive summaries."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=500,
                temperature=0.3
            )
            return response.choices[0].message.content.strip()

    except Exception as e:
        # Return original summary if LLM fails
        st.warning(f"LLM enhancement failed: {str(e)}")
        return summary

    return summary

def generate_pdf_report(data: Dict) -> bytes:
    """Generate comprehensive PDF report with benchmarks and risk analysis"""
    from reportlab.platypus import KeepTogether

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()

    # Get extracted data from session state if available
    extracted_data = st.session_state.get('extracted_data', None)

    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1f2937'),
        spaceAfter=30,
        alignment=1
    )

    story.append(Paragraph("DealGenie Pro - Investment Analysis", title_style))
    story.append(Spacer(1, 20))

    # Summary - use extracted_data if available
    if extracted_data:
        summary = generate_principal_summary(extracted_data)
    else:
        summary = generate_principal_summary(data)

    summary_style = ParagraphStyle(
        'SummaryStyle',
        parent=styles['Normal'],
        fontSize=11,
        leading=16,
        spaceAfter=12
    )

    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.HexColor('#1f2937'),
        spaceAfter=10,
        spaceBefore=15,
        fontName='Helvetica-Bold'
    )

    story.append(Paragraph("<b>INVESTMENT SUMMARY</b>", heading_style))
    # Preserve inline citations in summary
    formatted_summary = summary.replace('(', '<i>(').replace(')', ')</i>')
    story.append(Paragraph(formatted_summary, summary_style))
    story.append(Spacer(1, 20))

    # Key Metrics Table
    metrics_data = [
        ['Metric', 'Value'],
        ['Asset Class', data.get('asset_class', 'N/A')],
        ['Purchase Price', f"${data.get('purchase_price', 0):,.0f}"],
        ['NOI', f"${data.get('noi', 0):,.0f}"],
        ['Cap Rate', f"{(data.get('noi', 0) / max(data.get('purchase_price', 1), 1)) * 100:.2f}%"],
        ['Loan Amount', f"${data.get('loan_amount', 0):,.0f}"],
        ['Interest Rate', f"{data.get('interest_rate', 0) * 100:.2f}%"],
    ]

    metrics_table = Table(metrics_data, colWidths=[3*inch, 2*inch])
    metrics_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#374151')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    story.append(metrics_table)
    story.append(Spacer(1, 20))

    # Asset-Specific Metrics Section
    if data.get('asset_class'):
        story.append(Paragraph(f"<b>{data['asset_class'].upper()} SPECIFIC METRICS</b>", styles['Heading2']))

        asset_specific_data = []

        if data['asset_class'] == 'Office' and 'gla_sf' in data:
            asset_specific_data = [
                ['Metric', 'Value'],
                ['Gross Leasable Area', f"{data.get('gla_sf', 0):,.0f} SF"],
                ['WALT', f"{data.get('walt_years', 0):.1f} years"],
                ['TI - New Leases', f"${data.get('ti_new_psf', 0):.0f}/SF"],
                ['TI - Renewals', f"${data.get('ti_renewal_psf', 0):.0f}/SF"],
                ['LC - New', f"{data.get('lc_new_pct', 0)*100:.1f}%"],
                ['LC - Renewals', f"{data.get('lc_renewal_pct', 0)*100:.1f}%"],
                ['Tenant Count', f"{data.get('tenant_count', 0)}"],
                ['Top 5 Tenants', f"{data.get('top5_tenants_pct', 0)*100:.0f}%"],
            ]
        elif data['asset_class'] == 'Multifamily' and 'units' in data:
            asset_specific_data = [
                ['Metric', 'Value'],
                ['Total Units', f"{data.get('units', 0):,}"],
                ['Average Rent', f"${data.get('avg_rent', 0):,.0f}/month"],
                ['Market Rent', f"${data.get('market_rent', 0):,.0f}/month"],
                ['Occupancy', f"{data.get('occupancy_pct', 0)*100:.1f}%"],
                ['Expense Ratio', f"{data.get('expense_ratio', 0)*100:.1f}%"],
                ['Concessions', f"{data.get('concessions_months', 0):.1f} months"],
            ]
        elif data['asset_class'] == 'Industrial' and 'building_sf' in data:
            asset_specific_data = [
                ['Metric', 'Value'],
                ['Building Size', f"{data.get('building_sf', 0):,.0f} SF"],
                ['Clear Height', f"{data.get('clear_height_ft', 0)} ft"],
                ['Dock Doors', f"{data.get('dock_doors', 0)}"],
                ['Office Finish', f"{data.get('office_finish_pct', 0)*100:.1f}%"],
                ['Cold Storage', 'Yes' if data.get('cold_storage', False) else 'No'],
            ]
        elif data['asset_class'] == 'Retail' and 'gla_sf' in data:
            asset_specific_data = [
                ['Metric', 'Value'],
                ['Gross Leasable Area', f"{data.get('gla_sf', 0):,.0f} SF"],
                ['Anchor Tenant', data.get('anchor_tenant', 'N/A')],
                ['Anchor Term Remaining', f"{data.get('anchor_term_years', 0):.1f} years"],
                ['Co-Tenancy Clause', 'Yes' if data.get('co_tenancy_clause', False) else 'No'],
                ['Sales PSF', f"${data.get('sales_psf', 0):.0f}"],
            ]
        elif data['asset_class'] in ['Hotel', 'Hospitality'] and 'keys' in data:
            asset_specific_data = [
                ['Metric', 'Value'],
                ['Keys/Rooms', f"{data.get('keys', 0):,}"],
                ['ADR', f"${data.get('adr', 0):.0f}"],
                ['Occupancy', f"{data.get('occupancy_pct', 0)*100:.1f}%"],
                ['RevPAR', f"${data.get('revpar', 0):.0f}"],
                ['GOP Margin', f"{data.get('gop_margin_pct', 0)*100:.1f}%"],
                ['Brand/Flag', data.get('brand_flag', 'N/A')],
                ['PIP Cost per Key', f"${data.get('pip_cost_per_key', 0):,.0f}"],
            ]

        if asset_specific_data:
            asset_table = Table(asset_specific_data, colWidths=[3*inch, 2*inch])
            asset_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#374151')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ]))
            story.append(asset_table)
            story.append(Spacer(1, 20))

    # Benchmark Comparison Section
    if extracted_data and 'bench_compare' in extracted_data and extracted_data['bench_compare']:
        story.append(Paragraph("<b>BENCHMARK COMPARISON</b>", heading_style))

        bench_data = [['Metric', 'Your Value', 'Target Range', 'Status', 'Source']]

        for comp in extracted_data['bench_compare']:
            status = comp.get('status', 'OK')
            status_display = '✓ OK' if status == 'OK' else '⚠ Borderline' if status in ['Borderline', 'BORDERLINE'] else '✗ Offside'

            bench_data.append([
                comp.get('metric', 'Unknown'),
                str(comp.get('your_value', 'N/A')),
                comp.get('target_range', 'N/A'),
                status_display,
                comp.get('source', 'Market Data')
            ])

        bench_table = Table(bench_data, colWidths=[1.5*inch, 1.2*inch, 1.5*inch, 1*inch, 1.8*inch])
        bench_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#374151')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('FONTSIZE', (0, 1), (-1, -1), 9),
            ('GRID', (0, 0), (-1, -1), 1, colors.grey),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            # Color coding for status column
            ('TEXTCOLOR', (3, 1), (3, -1), colors.black),
        ]))

        # Apply conditional formatting to status column
        for i, row in enumerate(bench_data[1:], start=1):
            if '✓' in row[3]:
                bench_table.setStyle(TableStyle([('TEXTCOLOR', (3, i), (3, i), colors.green)]))
            elif '⚠' in row[3]:
                bench_table.setStyle(TableStyle([('TEXTCOLOR', (3, i), (3, i), colors.orange)]))
            elif '✗' in row[3]:
                bench_table.setStyle(TableStyle([('TEXTCOLOR', (3, i), (3, i), colors.red)]))

        story.append(bench_table)
        story.append(Spacer(1, 20))

    # Risk Analysis Section
    if extracted_data and 'risks_ranked' in extracted_data and extracted_data['risks_ranked']:
        story.append(Paragraph("<b>RISK ANALYSIS</b>", heading_style))

        # Filter for HIGH severity risks
        high_risks = [risk for risk in extracted_data['risks_ranked'] if risk.get('severity') == 'HIGH']

        if high_risks:
            for risk in high_risks:
                # Risk header with severity badge
                risk_text = f"<b>{risk.get('metric', 'Risk')}</b> - <font color='red'>HIGH SEVERITY</font><br/>"
                risk_text += f"<b>Current:</b> {risk.get('current_value', 'N/A')} | "
                risk_text += f"<b>Target:</b> {risk.get('target_value', 'N/A')}<br/>"

                if risk.get('explanation'):
                    risk_text += f"{risk['explanation']}<br/>"

                story.append(Paragraph(risk_text, summary_style))

                # Mitigations
                if 'mitigations' in risk and risk['mitigations']:
                    story.append(Paragraph("<b>Recommended Mitigations:</b>", summary_style))

                    mitigation_items = []
                    for i, mitigation in enumerate(risk['mitigations'], 1):
                        mit_text = f"{i}. {mitigation.get('action', 'Action')}"
                        if mitigation.get('dollar_impact'):
                            mit_text += f" <i>(Impact: ${mitigation['dollar_impact']:,})</i>"
                        mitigation_items.append(Paragraph(mit_text, summary_style))

                    for item in mitigation_items:
                        story.append(item)

                story.append(Spacer(1, 15))
        else:
            story.append(Paragraph("No high severity risks identified.", summary_style))
            story.append(Spacer(1, 15))

    story.append(PageBreak())

    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def generate_excel_export(data: Dict) -> bytes:
    """Generate Excel export with all data"""
    output = io.BytesIO()

    # Create Excel writer
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        # Summary sheet
        summary_data = {
            'Metric': ['Purchase Price', 'NOI', 'Cap Rate', 'Loan Amount', 'Interest Rate'],
            'Value': [
                f"${data.get('purchase_price', 0):,.0f}",
                f"${data.get('noi', 0):,.0f}",
                f"{(data.get('noi', 0) / max(data.get('purchase_price', 1), 1)) * 100:.2f}%",
                f"${data.get('loan_amount', 0):,.0f}",
                f"{data.get('interest_rate', 0) * 100:.2f}%"
            ]
        }
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Summary', index=False)

        # Input data sheet
        input_data = [(k, v) for k, v in data.items() if v is not None]
        if input_data:
            input_df = pd.DataFrame(input_data, columns=['Field', 'Value'])
            input_df.to_excel(writer, sheet_name='Input Data', index=False)

    output.seek(0)
    return output.getvalue()

def generate_chart_export(data: Dict) -> bytes:
    """Generate chart image for export"""
    fig, axes = plt.subplots(2, 2, figsize=(12, 10))
    fig.suptitle('DealGenie Investment Analysis', fontsize=16, fontweight='bold')

    # Calculate metrics
    cap_rate = (data.get('noi', 0) / max(data.get('purchase_price', 1), 1)) * 100
    ltv = (data.get('loan_amount', 0) / max(data.get('purchase_price', 1), 1)) * 100
    dscr = calculate_dscr(
        data.get('noi', 0),
        data.get('loan_amount', 0),
        data.get('interest_rate', 0.065),
        data.get('amort_years', 30)
    )

    # Metrics chart
    ax1 = axes[0, 0]
    metrics_names = ['Cap Rate', 'LTV', 'DSCR']
    metrics_values = [cap_rate/10, ltv/100, dscr]
    colors = ['#22c55e' if v > 0.5 else '#ef4444' for v in metrics_values]
    ax1.bar(metrics_names, metrics_values, color=colors)
    ax1.set_title('Key Metrics')
    ax1.set_ylim(0, 2)

    # Sensitivity matrix
    ax2 = axes[0, 1]
    ax2.set_title('Cap Rate Sensitivity')
    ax2.axis('off')

    # Cash flow projection
    ax3 = axes[1, 0]
    years = list(range(1, 6))
    cash_flows = [data.get('noi', 0) * (1.03 ** y) for y in years]
    ax3.plot(years, cash_flows, marker='o', color='purple', linewidth=2)
    ax3.set_title('NOI Projection')
    ax3.set_xlabel('Year')
    ax3.set_ylabel('NOI ($)')
    ax3.grid(True, alpha=0.3)

    # Value composition
    ax4 = axes[1, 1]
    sizes = [data.get('loan_amount', 0), max(data.get('purchase_price', 0) - data.get('loan_amount', 0), 0)]
    labels = ['Debt', 'Equity']
    ax4.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90, colors=['#ef4444', '#22c55e'])
    ax4.set_title('Capital Structure')

    # Save to bytes
    buffer = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    buffer.seek(0)
    plt.close()

    return buffer.getvalue()

def render_analysis(data: Dict):
    """Render analysis results with principal summary at top"""
    if not data or "purchase_price" not in data:
        st.info("Enter deal information to see analysis")
        return

    # Check if we have enhanced extracted data from CRE extraction engine
    extracted_data = st.session_state.get('extracted_data', None)

    # Generate principal summary with extracted data if available
    if extracted_data:
        summary = generate_principal_summary(extracted_data)
    else:
        # Fallback to legacy format
        summary = generate_principal_summary(data)

    # Calculate metrics for LLM context
    metrics = calculate_metrics_for_llm(data)

    # Use the new LLM-enhanced summary renderer
    render_summary_with_llm_option(summary, metrics)

    # Calculate metrics
    cap_rate = (data.get("noi", 0) / data.get("purchase_price", 1)) * 100
    ltv = (data.get("loan_amount", 0) / data.get("purchase_price", 1)) * 100
    dscr = calculate_dscr(
        data.get("noi", 0),
        data.get("loan_amount", 0),
        data.get("interest_rate", 0.065),
        data.get("amort_years", 30)
    )

    # Display asset-specific metrics if available
    if data.get("asset_class"):
        st.subheader(f"🏗️ {data['asset_class']} Specific Metrics")

        if data["asset_class"] == "Office" and 'gla_sf' in data:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("GLA", f"{data.get('gla_sf', 0):,.0f} SF",
                         help=create_metric_help_text("square_feet"))
            with col2:
                st.metric("WALT", f"{data.get('walt_years', 0):.1f} years",
                         help=create_metric_help_text("walt"))
            with col3:
                st.metric("TI New", f"${data.get('ti_new_psf', 0):.0f}/SF",
                         help=create_metric_help_text("tenant_improvement"))
            with col4:
                st.metric("Top 5 Tenants", f"{data.get('top5_tenants_pct', 0)*100:.0f}%",
                         help="Concentration risk from largest tenants")

        elif data["asset_class"] == "Multifamily" and 'units' in data:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Units", f"{data.get('units', 0):,}",
                         help=create_metric_help_text("units"))
            with col2:
                st.metric("Avg Rent", f"${data.get('avg_rent', 0):,.0f}",
                         help=create_metric_help_text("avg_rent"))
            with col3:
                st.metric("Occupancy", f"{data.get('occupancy_pct', 0)*100:.1f}%",
                         help=create_metric_help_text("occupancy"))
            with col4:
                st.metric("Expense Ratio", f"{data.get('expense_ratio', 0)*100:.1f}%",
                         help=create_metric_help_text("expense_ratio"))

        elif data["asset_class"] == "Industrial" and 'building_sf' in data:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Building SF", f"{data.get('building_sf', 0):,.0f}",
                         help=create_metric_help_text("square_feet"))
            with col2:
                st.metric("Clear Height", f"{data.get('clear_height_ft', 0)} ft",
                         help=create_metric_help_text("clear_height"))
            with col3:
                st.metric("Dock Doors", f"{data.get('dock_doors', 0)}",
                         help="Loading dock positions for truck access")
            with col4:
                st.metric("Office %", f"{data.get('office_finish_pct', 0)*100:.1f}%",
                         help="Percentage of space finished as office")

        elif data["asset_class"] == "Retail" and 'gla_sf' in data:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("GLA", f"{data.get('gla_sf', 0):,.0f} SF",
                         help=create_metric_help_text("square_feet"))
            with col2:
                st.metric("Anchor", data.get('anchor_tenant', 'N/A'),
                         help=create_metric_help_text("anchor_tenant"))
            with col3:
                st.metric("Anchor Term", f"{data.get('anchor_term_years', 0):.1f} yrs",
                         help="Remaining lease term for anchor tenant")
            with col4:
                st.metric("Sales/SF", f"${data.get('sales_psf', 0):.0f}",
                         help=create_metric_help_text("sales_psf"))

        elif data["asset_class"] in ["Hotel", "Hospitality"] and 'keys' in data:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Keys", f"{data.get('keys', 0):,}",
                         help=create_metric_help_text("keys"))
            with col2:
                st.metric("ADR", f"${data.get('adr', 0):.0f}",
                         help=create_metric_help_text("adr"))
            with col3:
                st.metric("RevPAR", f"${data.get('revpar', 0):.0f}",
                         help=create_metric_help_text("revpar"))
            with col4:
                st.metric("GOP Margin", f"{data.get('gop_margin_pct', 0)*100:.1f}%",
                         help=create_metric_help_text("gop_margin"))

    # Check if we have enhanced extracted data for tabbed interface
    if extracted_data:
        # Enhanced six-tab interface when extraction data is available
        tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
            "📊 Key Metrics", "✅ What We Know", "❌ What We Don't",
            "⚠️ Risk Analysis", "📈 Sensitivities", "🔍 Benchmarks"
        ])

        with tab1:
            # Tab 1: Key Metrics
            st.subheader("📈 Deal Metrics")

            # Main metrics cards
            col1, col2, col3, col4 = st.columns(4)

            with col1:
                st.markdown(f"""
                <div class="metric-card">
                    <h3 style="color: #667eea; margin: 0;">Cap Rate</h3>
                    <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{cap_rate:.2f}%</p>
                </div>
                """, unsafe_allow_html=True)

            with col2:
                st.markdown(f"""
                <div class="metric-card">
                    <h3 style="color: #667eea; margin: 0;">DSCR</h3>
                    <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{dscr:.2f}×</p>
                </div>
                """, unsafe_allow_html=True)

            with col3:
                st.markdown(f"""
                <div class="metric-card">
                    <h3 style="color: #667eea; margin: 0;">LTV</h3>
                    <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{ltv:.1f}%</p>
                </div>
                """, unsafe_allow_html=True)

            with col4:
                equity = data.get("purchase_price", 0) - data.get("loan_amount", 0)
                cash_on_cash = ((data.get("noi", 0) - data.get("loan_amount", 0) * data.get("interest_rate", 0.065)) / equity * 100) if equity > 0 else 0

                st.markdown(f"""
                <div class="metric-card">
                    <h3 style="color: #667eea; margin: 0;">Cash-on-Cash</h3>
                    <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{cash_on_cash:.1f}%</p>
                </div>
                """, unsafe_allow_html=True)

            # Analysis table
            st.markdown("---")
            st.subheader("📋 Full Analysis")

            analysis_data = {
                "Metric": ["Purchase Price", "NOI", "Cap Rate", "Loan Amount", "LTV",
                          "Interest Rate", "DSCR", "Equity Required", "Cash-on-Cash Return"],
                "Value": [
                    f"${data.get('purchase_price', 0):,.0f}",
                    f"${data.get('noi', 0):,.0f}",
                    f"{cap_rate:.2f}%",
                    f"${data.get('loan_amount', 0):,.0f}",
                    f"{ltv:.1f}%",
                    f"{data.get('interest_rate', 0.065)*100:.2f}%",
                    f"{dscr:.2f}x",
                    f"${equity:,.0f}",
                    f"{cash_on_cash:.1f}%"
                ]
            }

            df_analysis = pd.DataFrame(analysis_data)
            st.dataframe(df_analysis, use_container_width=True, hide_index=True)

        with tab2:
            # Tab 2: What We Know
            st.subheader("✅ Confirmed Data Points")

            # Get asset class and subclass for filtering
            asset_class = extracted_data.get('asset_class', 'office').lower()
            subclass = extracted_data.get('subclass', '').lower() or 'general'

            # Get filtered metrics for this asset class
            relevant_metrics = get_all_metrics_for_asset_class(asset_class, subclass)
            all_relevant_metrics = set()
            for category_metrics in relevant_metrics.values():
                all_relevant_metrics.update(category_metrics)

            # Check for new field_confidence structure
            if 'field_confidence' in extracted_data and extracted_data['field_confidence']:
                known_items = []

                # Process ingested fields with confidence
                ingested = extracted_data.get('ingested', {})
                field_confidence = extracted_data['field_confidence']

                for field_name, value in ingested.items():
                    # Skip fields not relevant to this asset class
                    if field_name.lower() not in all_relevant_metrics:
                        continue

                    if field_name in field_confidence:
                        conf_info = field_confidence[field_name]
                        confidence_level = conf_info.get('level', 'Medium')
                        reason = conf_info.get('reason', 'Found in document')

                        # Format value for display
                        if isinstance(value, float):
                            if field_name.endswith('_pct') or field_name.endswith('_rate'):
                                formatted_value = f"{value:.2%}"
                            elif value > 1000:
                                formatted_value = f"${value:,.0f}"
                            else:
                                formatted_value = f"{value:.3f}"
                        else:
                            formatted_value = str(value)

                        # Set badge color and icon based on confidence
                        if confidence_level == 'High':
                            badge_color = '#10b981'
                            icon = '🟢'
                        elif confidence_level == 'Medium':
                            badge_color = '#f59e0b'
                            icon = '🟡'
                        else:
                            badge_color = '#ef4444'
                            icon = '🔴'

                        known_items.append({
                            "Field": field_name.replace('_', ' ').title(),
                            "Value": formatted_value,
                            "Confidence": f"""<span style='background-color: {badge_color};
                                             color: white; padding: 3px 10px; border-radius: 12px;
                                             font-size: 11px; font-weight: 500;'>{icon} {confidence_level}</span>""",
                            "Source": f"<span style='color: #6b7280; font-size: 11px;'>{reason}</span>"
                        })

                # Also process derived fields
                derived = extracted_data.get('derived', {})
                for field_name, value in derived.items():
                    # Skip fields not relevant to this asset class
                    if field_name.lower() not in all_relevant_metrics:
                        continue

                    if not field_name.endswith('_calc') and field_name in field_confidence:
                        conf_info = field_confidence[field_name]
                        confidence_level = conf_info.get('level', 'Low')
                        reason = conf_info.get('reason', 'Calculated')

                        # Format value
                        if isinstance(value, float):
                            if field_name in ['cap_rate', 'dscr', 'ltv', 'debt_yield']:
                                formatted_value = f"{value:.3f}"
                            elif field_name in ['exit_value', 'net_sale_proceeds']:
                                formatted_value = f"${value:,.0f}"
                            else:
                                formatted_value = f"{value:.3f}"
                        else:
                            formatted_value = str(value)

                        # Badge styling for calculated fields
                        badge_color = '#9333ea'  # Purple for calculated
                        icon = '🔮'

                        known_items.append({
                            "Field": field_name.replace('_', ' ').title(),
                            "Value": formatted_value,
                            "Confidence": f"""<span style='background-color: {badge_color};
                                             color: white; padding: 3px 10px; border-radius: 12px;
                                             font-size: 11px; font-weight: 500;'>{icon} Calculated</span>""",
                            "Source": f"<span style='color: #6b7280; font-size: 11px;'>{reason}</span>"
                        })

                if known_items:
                    df_known = pd.DataFrame(known_items)
                    # Display as HTML table with enhanced styling
                    table_html = df_known.to_html(escape=False, index=False)
                    st.markdown(f"""
                    <style>
                    table {{
                        width: 100%;
                        border-collapse: collapse;
                    }}
                    th {{
                        background-color: #f3f4f6;
                        padding: 12px;
                        text-align: left;
                        font-weight: 600;
                    }}
                    td {{
                        padding: 10px 12px;
                        border-bottom: 1px solid #e5e7eb;
                    }}
                    tr:hover {{
                        background-color: #f9fafb;
                    }}
                    </style>
                    {table_html}
                    """, unsafe_allow_html=True)

                    # Add legend
                    st.caption("""
                    **Confidence Levels:**
                    🟢 **High** - Found in table or with explicit label |
                    🟡 **Medium** - Pattern matched in text |
                    🔴 **Low** - Inferred or using defaults |
                    🔮 **Calculated** - Derived from other fields
                    """)
                else:
                    st.info("No data points extracted yet.")

            # Fallback to old structure if field_confidence not available
            elif 'known' in extracted_data and extracted_data['known']:
                known_items = []
                for item in extracted_data['known']:
                    known_items.append({
                        "Field": item if isinstance(item, str) else item.get('field', 'Unknown'),
                        "Value": item.get('value', 'N/A') if isinstance(item, dict) else 'N/A',
                        "Confidence": "Medium"
                    })
                df_known = pd.DataFrame(known_items)
                st.dataframe(df_known, use_container_width=True, hide_index=True)
            else:
                st.info("No extracted data points available. Upload a document or enter data manually.")

        with tab3:
            # Tab 3: What We Don't Know
            st.subheader("❌ Missing Data Points")

            if 'unknown' in extracted_data and extracted_data['unknown']:
                # Group unknown items by type
                derived_metrics = []
                primary_fields = []

                for item in extracted_data['unknown']:
                    if 'missing' in item and item['missing']:
                        derived_metrics.append(item)
                    else:
                        primary_fields.append(item)

                # Display derived metrics that cannot be calculated
                if derived_metrics:
                    st.markdown("### 📊 **Metrics That Cannot Be Calculated**")
                    for item in derived_metrics:
                        metric_name = item.get('metric', 'Unknown').replace('_', ' ').title()
                        missing_fields = item.get('missing', [])
                        explanation = item.get('because', '')

                        # Create expandable card for each metric
                        with st.expander(f"🔍 **{metric_name}**", expanded=False):
                            # Show explanation
                            st.info(f"💡 {explanation}")

                            # Show missing fields
                            if missing_fields:
                                st.write("**Missing Required Fields:**")
                                cols = st.columns(2)
                                for i, field in enumerate(missing_fields):
                                    with cols[i % 2]:
                                        field_display = field.replace('_', ' ').title()
                                        st.markdown(f"""
                                        <div style='background: #fee2e2; padding: 8px 12px;
                                                  border-radius: 8px; margin: 4px 0;
                                                  border-left: 3px solid #ef4444;'>
                                            <span style='color: #991b1b; font-weight: 500;'>
                                                ❌ {field_display}
                                            </span>
                                        </div>
                                        """, unsafe_allow_html=True)

                            # Show calculation formula if available
                            formula_map = {
                                "irr": "IRR = Solve for r where: NPV = 0 = -Equity + Σ(CFt/(1+r)^t)",
                                "equity_multiple": "Equity Multiple = Total Distributions ÷ Initial Equity",
                                "exit_value": "Exit Value = NOI(Year N) × (1 + Growth)^N ÷ Exit Cap",
                                "dscr": "DSCR = Net Operating Income ÷ Annual Debt Service",
                                "ltv": "LTV = Loan Amount ÷ Purchase Price",
                                "cap_rate": "Cap Rate = NOI ÷ Purchase Price"
                            }

                            if item['metric'] in formula_map:
                                st.code(formula_map[item['metric']], language='text')

                # Display missing primary fields
                if primary_fields:
                    st.markdown("---")
                    st.markdown("### 📝 **Missing Primary Data Fields**")

                    # Group by importance
                    critical_fields = []
                    optional_fields = []

                    for item in primary_fields:
                        if any(keyword in item.get('because', '').lower()
                               for keyword in ['critical', 'required', 'essential']):
                            critical_fields.append(item)
                        else:
                            optional_fields.append(item)

                    if critical_fields:
                        st.markdown("**🔴 Critical Fields:**")
                        for item in critical_fields:
                            field_name = item.get('metric', '').replace('_', ' ').title()
                            reason = item.get('because', item.get('description', ''))
                            st.markdown(f"""
                            <div style='background: #fef2f2; padding: 12px;
                                      border-radius: 8px; margin: 8px 0;
                                      border-left: 4px solid #dc2626;'>
                                <strong style='color: #7f1d1d;'>{field_name}</strong><br/>
                                <span style='color: #991b1b; font-size: 0.9em;'>{reason}</span>
                            </div>
                            """, unsafe_allow_html=True)

                    if optional_fields:
                        st.markdown("**🟡 Additional Fields:**")
                        for item in optional_fields:
                            field_name = item.get('metric', '').replace('_', ' ').title()
                            reason = item.get('because', item.get('description', ''))
                            st.markdown(f"""
                            <div style='background: #fffbeb; padding: 12px;
                                      border-radius: 8px; margin: 8px 0;
                                      border-left: 4px solid #f59e0b;'>
                                <strong style='color: #78350f;'>{field_name}</strong><br/>
                                <span style='color: #92400e; font-size: 0.9em;'>{reason}</span>
                            </div>
                            """, unsafe_allow_html=True)

                # Show completeness score
                if 'completeness' in extracted_data:
                    st.markdown("---")
                    completeness = extracted_data['completeness']
                    percent = completeness.get('percent', 0)
                    filled = completeness.get('filled', 0)
                    total = completeness.get('required', 0)

                    # Progress bar color based on percentage
                    if percent >= 80:
                        color = '#10b981'
                    elif percent >= 60:
                        color = '#f59e0b'
                    else:
                        color = '#ef4444'

                    st.markdown(f"""
                    <div style='background: #f3f4f6; padding: 1rem; border-radius: 8px;'>
                        <h4 style='margin: 0 0 0.5rem 0;'>📊 Data Completeness</h4>
                        <div style='background: #e5e7eb; border-radius: 20px; height: 30px; position: relative;'>
                            <div style='background: {color}; border-radius: 20px; height: 100%;
                                      width: {percent}%; transition: width 0.5s ease;'></div>
                            <span style='position: absolute; top: 50%; left: 50%;
                                       transform: translate(-50%, -50%); font-weight: bold;'>
                                {percent:.1f}% ({filled}/{total} fields)
                            </span>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)

            else:
                st.success("✅ All critical data points have been captured!")

        with tab4:
            # Tab 4: Risk Analysis
            st.subheader("⚠️ Risk Assessment")

            # Get asset class and subclass for filtering
            asset_class = extracted_data.get('asset_class', 'office').lower()
            subclass = extracted_data.get('subclass', '').lower() or 'general'

            # Get filtered metrics for this asset class
            relevant_metrics = get_all_metrics_for_asset_class(asset_class, subclass)
            all_relevant_metrics = set()
            for category_metrics in relevant_metrics.values():
                all_relevant_metrics.update(category_metrics)

            if 'risks_ranked' in extracted_data and extracted_data['risks_ranked']:
                # Filter risks to only show relevant metrics
                relevant_risks = []
                for risk in extracted_data['risks_ranked']:
                    risk_metric = risk.get('metric', '').lower().replace(' ', '_')
                    if risk_metric in all_relevant_metrics:
                        relevant_risks.append(risk)

                for risk in relevant_risks:
                    severity = risk.get('severity', 'MEDIUM')

                    # Color coding for risk cards
                    if severity == 'HIGH':
                        card_color = 'linear-gradient(135deg, #ef4444 0%, #dc2626 100%)'
                        badge_color = '#991b1b'
                    elif severity == 'MEDIUM':
                        card_color = 'linear-gradient(135deg, #f59e0b 0%, #d97706 100%)'
                        badge_color = '#92400e'
                    else:
                        card_color = 'linear-gradient(135deg, #10b981 0%, #059669 100%)'
                        badge_color = '#064e3b'

                    st.markdown(f"""
                    <div style="background: {card_color}; padding: 1.5rem; border-radius: 12px;
                                color: white; margin-bottom: 1rem;">
                        <div style="display: flex; justify-content: space-between; align-items: center;">
                            <h3 style="margin: 0;">{risk.get('metric', 'Risk')}</h3>
                            <span style="background: {badge_color}; padding: 4px 12px;
                                       border-radius: 20px; font-weight: bold;">{severity}</span>
                        </div>
                        <p style="margin-top: 1rem; margin-bottom: 0.5rem;">
                            <strong>Current:</strong> {risk.get('current_value', 'N/A')} |
                            <strong>Target:</strong> {risk.get('target_value', 'N/A')}
                        </p>
                        <p style="margin-top: 0.5rem;">{risk.get('explanation', '')}</p>

                        <h4 style="margin-top: 1rem; margin-bottom: 0.5rem;">Mitigations:</h4>
                        <ol style="margin: 0; padding-left: 1.5rem;">
                    """, unsafe_allow_html=True)

                    # Add mitigations if available
                    if 'mitigations' in risk:
                        for i, mitigation in enumerate(risk['mitigations'], 1):
                            impact = mitigation.get('dollar_impact', '')
                            st.markdown(f"""
                            <li>{mitigation.get('action', 'Action')}
                                {f'<em>(Impact: ${impact:,})</em>' if impact else ''}</li>
                            """, unsafe_allow_html=True)

                    st.markdown("</ol></div>", unsafe_allow_html=True)

                if not relevant_risks:
                    st.info("No risks identified for this asset class.")
            else:
                st.info("No risk analysis available. Complete data entry for risk assessment.")

        with tab5:
            # Tab 5: Sensitivities
            st.subheader("📈 Sensitivity Analysis")

            # Exit cap sensitivity
            st.markdown("#### Exit Cap Rate Sensitivity")

            # Create sliders for assumptions
            col1, col2 = st.columns(2)
            with col1:
                base_exit_cap = st.slider(
                    "Base Exit Cap Rate (%)",
                    min_value=3.0,
                    max_value=12.0,
                    value=cap_rate + 0.5,
                    step=0.25,
                    help="Assumed cap rate at sale"
                )

            with col2:
                hold_period = st.slider(
                    "Hold Period (Years)",
                    min_value=1,
                    max_value=10,
                    value=5,
                    help="Investment hold period"
                )

            # Calculate sensitivity scenarios
            scenarios = []
            exit_caps = [base_exit_cap - 0.5, base_exit_cap, base_exit_cap + 0.5, base_exit_cap + 1.0]

            for exit_cap_scenario in exit_caps:
                future_noi = data.get("noi", 0) * (1.03 ** hold_period)
                exit_value = future_noi / (exit_cap_scenario / 100) if exit_cap_scenario > 0 else 0

                # Assume 10% loan paydown
                remaining_balance = data.get("loan_amount", 0) * 0.9
                net_proceeds = exit_value - remaining_balance

                equity = data.get("purchase_price", 0) - data.get("loan_amount", 0)
                equity_multiple = net_proceeds / equity if equity > 0 else 0

                scenarios.append({
                    'Exit Cap': f"{exit_cap_scenario:.2f}%",
                    'Exit Value': exit_value,
                    'Equity Multiple': equity_multiple
                })

            # Create interactive chart
            df_scenarios = pd.DataFrame(scenarios)

            fig = go.Figure()

            # Add exit value line
            fig.add_trace(go.Scatter(
                x=df_scenarios['Exit Cap'],
                y=df_scenarios['Exit Value'],
                name='Exit Value',
                line=dict(color='#667eea', width=3),
                mode='lines+markers'
            ))

            # Add equity multiple on secondary y-axis
            fig.add_trace(go.Scatter(
                x=df_scenarios['Exit Cap'],
                y=df_scenarios['Equity Multiple'],
                name='Equity Multiple',
                line=dict(color='#10b981', width=3),
                mode='lines+markers',
                yaxis='y2'
            ))

            fig.update_layout(
                title="Exit Cap Rate Impact on Returns",
                xaxis_title="Exit Cap Rate",
                yaxis_title="Exit Value ($)",
                yaxis2=dict(
                    title="Equity Multiple (x)",
                    overlaying='y',
                    side='right'
                ),
                hovermode='x unified',
                height=400,
                showlegend=True
            )

            st.plotly_chart(fig, use_container_width=True)

            # Display sensitivity table
            st.markdown("#### Sensitivity Table")
            sensitivity_df = pd.DataFrame({
                "Exit Cap": df_scenarios['Exit Cap'],
                "Exit Value": [f"${v:,.0f}" for v in df_scenarios['Exit Value']],
                "Equity Multiple": [f"{em:.2f}x" for em in df_scenarios['Equity Multiple']]
            })
            st.dataframe(sensitivity_df, use_container_width=True, hide_index=True)

        with tab6:
            # Tab 6: Benchmarks
            st.subheader("🔍 Benchmark Comparison")

            if 'bench_compare' in extracted_data and extracted_data['bench_compare']:
                # Get asset class and subclass from extracted data
                asset_class = extracted_data.get('asset_class', 'office').lower()
                subclass = extracted_data.get('subclass', '').lower() or 'general'

                # Get filtered metrics for this asset class
                relevant_metrics = get_all_metrics_for_asset_class(asset_class, subclass)
                all_relevant_metrics = set()
                for category_metrics in relevant_metrics.values():
                    all_relevant_metrics.update(category_metrics)

                benchmark_data = []

                for comp in extracted_data['bench_compare']:
                    # Only include metrics that are relevant to this asset class
                    metric_name = comp.get('metric', '').lower().replace(' ', '_')
                    if metric_name not in all_relevant_metrics:
                        continue  # Skip irrelevant metrics

                    # Determine status
                    status = comp.get('status', 'OK')
                    if status == 'OK':
                        status_color = '#10b981'
                        status_text = '✅ OK'
                    elif status in ['Borderline', 'BORDERLINE']:
                        status_color = '#f59e0b'
                        status_text = '⚠️ Borderline'
                    else:
                        status_color = '#ef4444'
                        status_text = '❌ Offside'

                    benchmark_data.append({
                        "Metric": comp.get('metric', 'Unknown'),
                        "Your Value": comp.get('your_value', 'N/A'),
                        "Target Range": comp.get('target_range', 'N/A'),
                        "Status_HTML": f"<span style='color: {status_color}; font-weight: bold;'>{status_text}</span>",
                        "Source": comp.get('source', 'Market Data')
                    })

                if benchmark_data:
                    df_bench = pd.DataFrame(benchmark_data)
                    # Display as HTML table with styled status column
                    st.markdown(df_bench.to_html(escape=False, index=False), unsafe_allow_html=True)
                else:
                    st.info("No relevant benchmark comparisons for this asset class.")
            else:
                st.info("Benchmark comparison will be available once deal metrics are calculated.")

    else:
        # Legacy display without extraction data
        st.subheader("📈 Key Metrics")
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.markdown(f"""
            <div class="metric-card">
                <h3 style="color: #667eea; margin: 0;">Cap Rate</h3>
                <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{cap_rate:.2f}%</p>
            </div>
            """, unsafe_allow_html=True)

        with col2:
            st.markdown(f"""
            <div class="metric-card">
                <h3 style="color: #667eea; margin: 0;">DSCR</h3>
                <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{dscr:.2f}×</p>
            </div>
            """, unsafe_allow_html=True)

        with col3:
            st.markdown(f"""
            <div class="metric-card">
                <h3 style="color: #667eea; margin: 0;">LTV</h3>
                <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{ltv:.1f}%</p>
            </div>
            """, unsafe_allow_html=True)

        with col4:
            equity = data.get("purchase_price", 0) - data.get("loan_amount", 0)
            cash_on_cash = ((data.get("noi", 0) - data.get("loan_amount", 0) * data.get("interest_rate", 0.065)) / equity * 100) if equity > 0 else 0

            st.markdown(f"""
            <div class="metric-card">
                <h3 style="color: #667eea; margin: 0;">Cash-on-Cash</h3>
                <p style="font-size: 2rem; font-weight: bold; margin: 0.5rem 0;">{cash_on_cash:.1f}%</p>
            </div>
            """, unsafe_allow_html=True)

        # Benchmark evaluation
        st.subheader("🎯 Benchmark Analysis")

        metrics = {
            "cap_rate": cap_rate,
            "dscr": dscr,
            "ltv": ltv
        }

        evaluations = evaluate_against_benchmarks(
            data.get("asset_class", "Office"),
            metrics
        )

        for eval in evaluations:
            status_class = f"status-{eval['status']}"
            st.markdown(f"""
            <div style="padding: 0.5rem; margin: 0.5rem 0;">
                <span class="status-badge {status_class}">{eval['status'].upper()}</span>
                <span style="margin-left: 1rem; font-weight: 600;">{eval['metric']}: {eval['value']:.2f}</span>
                <span style="margin-left: 1rem; color: #6b7280;">Benchmark: {eval['benchmark']}</span>
            </div>
            """, unsafe_allow_html=True)

        # Cash flow projection
        st.subheader("💰 5-Year Cash Flow Projection")

        years = list(range(1, 6))
        noi_growth = 1.03  # 3% annual growth
        cash_flows = []

        for year in years:
            year_noi = data.get("noi", 0) * (noi_growth ** (year - 1))
            debt_service = data.get("loan_amount", 0) * data.get("interest_rate", 0.065)
            cash_flow = year_noi - debt_service
            cash_flows.append(cash_flow)

        fig = go.Figure()
        fig.add_trace(go.Bar(
            x=years,
            y=cash_flows,
            marker_color='#667eea',
            text=[f'${cf/1000:.0f}K' for cf in cash_flows],
            textposition='outside'
        ))

        fig.update_layout(
            title="Annual Cash Flow After Debt Service",
            xaxis_title="Year",
            yaxis_title="Cash Flow ($)",
            height=400,
            showlegend=False,
            plot_bgcolor='rgba(0,0,0,0)',
            paper_bgcolor='rgba(0,0,0,0)'
        )

        st.plotly_chart(fig, use_container_width=True)

def render_profile_and_templates():
    """Render user profile and template management in sidebar"""
    with st.sidebar:
        # Active Template Indicator (at the top)
        st.markdown("### 📌 Active Template")

        if st.session_state.get('template_loaded', False):
            # Custom template is loaded
            template_name = st.session_state.get('active_template_name', 'Unknown')
            template_date = st.session_state.get('active_template_date', 'Unknown')

            # Format date if available
            if template_date != "Unknown":
                try:
                    formatted_date = datetime.fromisoformat(template_date).strftime("%Y-%m-%d %H:%M")
                except:
                    formatted_date = template_date
            else:
                formatted_date = "Unknown"

            st.success(f"**{template_name}**")
            st.caption(f"📅 Created: {formatted_date}")
        else:
            # Using built-in benchmarks
            st.info("**Built-in Benchmarks**")
            st.caption("Using default DealGenie benchmark library")

        st.markdown("---")
        st.markdown("### 👤 User Profile")

        # Profile Name
        if 'profile_name' not in st.session_state:
            st.session_state.profile_name = ""

        profile_name = st.text_input(
            "Profile Name",
            value=st.session_state.profile_name,
            placeholder="Enter your name or company",
            key="profile_name_input"
        )
        st.session_state.profile_name = profile_name

        st.markdown("---")
        st.markdown("### 📋 Template Management")

        # Get available templates
        available_templates = list_templates()
        template_options = ["Default"] + available_templates + ["New Template"]

        # Template selector
        if 'selected_template' not in st.session_state:
            st.session_state.selected_template = "Default"

        selected_template = st.selectbox(
            "Template",
            options=template_options,
            index=template_options.index(st.session_state.selected_template) if st.session_state.selected_template in template_options else 0,
            key="template_selector"
        )

        # If "New Template" is selected, show input for new template name
        new_template_name = None
        if selected_template == "New Template":
            new_template_name = st.text_input(
                "New Template Name",
                placeholder="Enter template name",
                key="new_template_name_input"
            )

        st.markdown("---")

        # Template action buttons
        col1, col2 = st.columns(2)

        with col1:
            if st.button("💾 Save", use_container_width=True, help="Save current settings as template"):
                # Determine template name
                if selected_template == "New Template":
                    template_to_save = new_template_name
                elif selected_template == "Default":
                    st.warning("Cannot overwrite Default template. Please select 'New Template'.")
                    template_to_save = None
                else:
                    template_to_save = selected_template

                if template_to_save and template_to_save.strip():
                    # Gather settings from session state
                    settings_dict = {
                        "benchmark_overrides": st.session_state.get("benchmark_overrides", {}),
                        "custom_dd_items": st.session_state.get("custom_dd_items", {}),
                        "profile_name": st.session_state.get("profile_name", "")
                    }

                    if save_template(template_to_save, settings_dict):
                        st.success(f"✅ Template '{template_to_save}' saved!")
                        st.session_state.selected_template = template_to_save
                        st.rerun()
                elif selected_template == "New Template":
                    st.warning("Please enter a template name")

        with col2:
            if st.button("📂 Load", use_container_width=True, help="Load template settings"):
                if selected_template == "Default":
                    # Reset to defaults
                    if 'benchmark_overrides' in st.session_state:
                        st.session_state.benchmark_overrides = {}
                    if 'custom_dd_items' in st.session_state:
                        st.session_state.custom_dd_items = {}
                    st.success("✅ Reset to default settings")
                    st.rerun()
                elif selected_template == "New Template":
                    st.warning("Please select an existing template to load")
                else:
                    if load_template(selected_template):
                        st.success(f"✅ Template '{selected_template}' loaded!")
                        st.session_state.selected_template = selected_template
                        st.rerun()

        # Delete template option (only for custom templates)
        if selected_template not in ["Default", "New Template"] and available_templates:
            st.markdown("---")
            if st.button("🗑️ Delete Template", use_container_width=True, type="secondary"):
                if delete_template(selected_template):
                    st.success(f"✅ Template '{selected_template}' deleted")
                    st.session_state.selected_template = "Default"
                    st.rerun()

        # Export/Import Section
        st.markdown("---")
        st.markdown("**📤 Export / 📥 Import**")

        # Export current settings as JSON
        export_data = export_template_to_json()
        export_json = json.dumps(export_data, indent=2)
        export_filename = f"{st.session_state.get('profile_name', 'template').replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.json"

        st.download_button(
            label="📥 Export Template as JSON",
            data=export_json,
            file_name=export_filename,
            mime="application/json",
            use_container_width=True,
            help="Download current settings as a JSON file"
        )

        # Import template from JSON file
        uploaded_file = st.file_uploader(
            "📤 Import Template",
            type=['json'],
            help="Upload a previously exported template JSON file",
            key="template_uploader"
        )

        if uploaded_file is not None:
            try:
                # Read and parse JSON
                json_content = uploaded_file.read().decode('utf-8')
                imported_data = json.loads(json_content)

                # Import and validate
                if import_template_from_json(imported_data):
                    st.success(f"✅ Template imported successfully!")
                    st.session_state.selected_template = "Default"  # Reset selector
                    st.rerun()

            except json.JSONDecodeError as e:
                st.error(f"❌ Invalid JSON file: {str(e)}")
            except Exception as e:
                st.error(f"❌ Error reading file: {str(e)}")

        # Show template metadata
        st.markdown("---")
        st.markdown("**📊 Template Metadata**")

        # Display metadata based on source
        if selected_template not in ["Default", "New Template"]:
            # Show saved template metadata
            try:
                template_path = TEMPLATES_DIR / f"{selected_template}.json"
                if template_path.exists():
                    with open(template_path, 'r') as f:
                        template_data = json.load(f)

                    created_date = template_data.get("created_date", "Unknown")
                    if created_date != "Unknown":
                        created_date = datetime.fromisoformat(created_date).strftime("%Y-%m-%d %H:%M")

                    # Count customizations
                    num_benchmark_overrides = 0
                    if "benchmark_overrides" in template_data:
                        for asset_class, subclasses in template_data["benchmark_overrides"].items():
                            for subclass, metrics in subclasses.items():
                                num_benchmark_overrides += len(metrics)

                    num_custom_dd = sum(len(items) for items in template_data.get("custom_dd_items", {}).values())

                    st.caption(f"**Created:** {created_date}")
                    st.caption(f"**Custom Benchmarks:** {num_benchmark_overrides}")
                    st.caption(f"**Custom DD Items:** {num_custom_dd}")

            except Exception as e:
                st.caption("Unable to load metadata")

        elif hasattr(st.session_state, 'imported_template_name'):
            # Show imported template metadata
            imported_date = st.session_state.get('imported_template_date', 'Unknown')
            if imported_date != "Unknown":
                try:
                    imported_date = datetime.fromisoformat(imported_date).strftime("%Y-%m-%d %H:%M")
                except:
                    pass

            # Count current customizations
            num_benchmark_overrides = 0
            if hasattr(st.session_state, 'benchmark_overrides'):
                for asset_class, subclasses in st.session_state.benchmark_overrides.items():
                    for subclass, metrics in subclasses.items():
                        num_benchmark_overrides += len(metrics)

            num_custom_dd = sum(len(items) for items in st.session_state.get("custom_dd_items", {}).values())

            st.caption(f"**Source:** Imported ({st.session_state.imported_template_name})")
            st.caption(f"**Created:** {imported_date}")
            st.caption(f"**Custom Benchmarks:** {num_benchmark_overrides}")
            st.caption(f"**Custom DD Items:** {num_custom_dd}")

        else:
            # Show current state metadata
            num_benchmark_overrides = 0
            if hasattr(st.session_state, 'benchmark_overrides'):
                for asset_class, subclasses in st.session_state.benchmark_overrides.items():
                    for subclass, metrics in subclasses.items():
                        num_benchmark_overrides += len(metrics)

            num_custom_dd = sum(len(items) for items in st.session_state.get("custom_dd_items", {}).values())

            st.caption(f"**Custom Benchmarks:** {num_benchmark_overrides}")
            st.caption(f"**Custom DD Items:** {num_custom_dd}")

        st.markdown("---")

def main():
    """Main application entry point"""
    inject_custom_css()
    render_header()

    # ============================================================================
    # TEMPLATE PERSISTENCE & INITIALIZATION
    # ============================================================================

    # Ensure templates directory exists
    os.makedirs('data/templates', exist_ok=True)

    # Initialize template tracking flags
    if 'template_loaded' not in st.session_state:
        st.session_state.template_loaded = False
        st.session_state.active_template_name = None
        st.session_state.active_template_date = None

    # Load default template on first run (persists across sessions)
    if 'default_template_loaded' not in st.session_state:
        # Check if default template file exists
        default_template = get_default_template()

        if default_template:
            # Attempt to load the default template
            try:
                if load_template(default_template):
                    st.session_state.selected_template = default_template
                    st.session_state.template_loaded = True
                    st.session_state.active_template_name = default_template

                    # Get template metadata for last modified date
                    metadata = get_template_metadata(default_template)
                    if metadata:
                        st.session_state.active_template_date = metadata.get("created_date", "Unknown")
                else:
                    # Default template failed to load, use built-in benchmarks
                    st.session_state.template_loaded = False
                    st.session_state.active_template_name = "Built-in Benchmarks"
                    st.session_state.active_template_date = None
            except Exception as e:
                # Error loading default template, fall back to built-in
                st.session_state.template_loaded = False
                st.session_state.active_template_name = "Built-in Benchmarks"
                st.session_state.active_template_date = None
        else:
            # No default template set, use built-in benchmark library
            st.session_state.template_loaded = False
            st.session_state.active_template_name = "Built-in Benchmarks"
            st.session_state.active_template_date = None

        st.session_state.default_template_loaded = True

    # Render profile and template management in sidebar
    render_profile_and_templates()

    # Render API settings in sidebar
    render_api_settings()

    # Create tabs
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "📊 Analysis",
        "📚 Benchmarks",
        "📋 Due Diligence",
        "📄 Reports",
        "⚙️ Settings"
    ])

    with tab1:
        data = render_input_section()
        if data:
            render_analysis(data)
            # Store data in session state for other tabs
            st.session_state.analysis_data = data

    with tab2:
        st.header("📚 Industry Benchmarks")
        st.caption("Source-attributed market data with institutional standards")

        # Asset class and subclass selection
        col1, col2 = st.columns(2)
        with col1:
            selected_asset = st.selectbox(
                "Select Asset Class",
                options=list(BENCHMARK_DATA.keys()),
                format_func=lambda x: x.replace("_", " ").title()
            )

        with col2:
            if selected_asset in BENCHMARK_DATA:
                selected_subclass = st.selectbox(
                    "Select Property Type",
                    options=list(BENCHMARK_DATA[selected_asset].keys()),
                    format_func=lambda x: x.replace("_", " ").title()
                )
            else:
                selected_subclass = None

        if selected_asset and selected_subclass:
            # Initialize benchmark overrides in session state
            if 'benchmark_overrides' not in st.session_state:
                st.session_state.benchmark_overrides = {}
            if selected_asset not in st.session_state.benchmark_overrides:
                st.session_state.benchmark_overrides[selected_asset] = {}
            if selected_subclass not in st.session_state.benchmark_overrides[selected_asset]:
                st.session_state.benchmark_overrides[selected_asset][selected_subclass] = {}

            # Get benchmarks - use overrides if available, otherwise defaults
            default_benchmarks = BENCHMARK_DATA[selected_asset][selected_subclass]
            overrides = st.session_state.benchmark_overrides[selected_asset][selected_subclass]

            # Merge defaults with overrides
            selected_benchmarks = {}
            for metric, values in default_benchmarks.items():
                if metric in overrides:
                    selected_benchmarks[metric] = overrides[metric]
                else:
                    selected_benchmarks[metric] = values

            # Add edit mode toggle and override management
            col1, col2, col3 = st.columns([2, 1, 1])
            with col1:
                edit_mode = st.checkbox("🔧 Edit Benchmarks", help="Enable editing of benchmark values and sources")

            with col2:
                # Show override count
                override_count = len(overrides)
                if override_count > 0:
                    st.info(f"📝 {override_count} custom overrides active")

            with col3:
                # Clear overrides button
                if override_count > 0:
                    if st.button("🗑️ Clear All Overrides"):
                        st.session_state.benchmark_overrides[selected_asset][selected_subclass] = {}
                        st.success("Cleared all benchmark overrides")
                        st.rerun()

            # Show active overrides if any
            if overrides and not edit_mode:
                with st.expander(f"📋 Active Custom Overrides ({override_count})", expanded=False):
                    override_data = []
                    for metric, values in overrides.items():
                        override_data.append({
                            "Metric": metric.replace("_", " ").title(),
                            "Min": values[0],
                            "Preferred": values[1],
                            "Max": values[2],
                            "Source": values[3] if len(values) > 3 else "Custom"
                        })
                    st.dataframe(pd.DataFrame(override_data), use_container_width=True, hide_index=True)

            # Get filtered metrics for this asset class and subclass
            metric_groups = get_all_metrics_for_asset_class(selected_asset, selected_subclass)

            # Display metrics in expandable groups
            for group_name, metric_list in metric_groups.items():
                # Filter to only metrics that exist in selected benchmarks
                available_metrics = [m for m in metric_list if m in selected_benchmarks]

                if available_metrics:
                    with st.expander(f"{group_name} ({len(available_metrics)} metrics)", expanded=("Leverage" in group_name)):

                        if edit_mode:
                            # EDIT MODE - Create editable dataframe
                            edit_data = []
                            metric_map = {}  # Map display names back to metric keys

                            for metric in available_metrics:
                                if metric in selected_benchmarks:
                                    bench_data = selected_benchmarks[metric]
                                    metric_info = get_metric_info(metric)
                                    unit = metric_info.get("unit", "")

                                    # Store raw values for editing
                                    display_name = metric.replace("_", " ").title()
                                    metric_map[display_name] = metric

                                    # Extract raw values (convert percentages back to decimals for consistency)
                                    if isinstance(bench_data, (list, tuple)):
                                        min_val = bench_data[0] if len(bench_data) > 0 else 0
                                        pref_val = bench_data[1] if len(bench_data) > 1 else 0
                                        max_val = bench_data[2] if len(bench_data) > 2 else 0
                                        source = bench_data[3] if len(bench_data) > 3 else "Industry Standard"
                                    else:
                                        min_val = pref_val = max_val = 0
                                        source = "Industry Standard"

                                    edit_data.append({
                                        "Metric": display_name,
                                        "Min": min_val if unit != "%" else min_val * 100,
                                        "Preferred": pref_val if unit != "%" else pref_val * 100,
                                        "Max": max_val if unit != "%" else max_val * 100,
                                        "Unit": unit,
                                        "Source": source
                                    })

                            if edit_data:
                                df_edit = pd.DataFrame(edit_data)

                                # Use data_editor for editing
                                edited_df = st.data_editor(
                                    df_edit,
                                    use_container_width=True,
                                    hide_index=True,
                                    num_rows="fixed",
                                    disabled=["Metric", "Unit"],
                                    column_config={
                                        "Metric": st.column_config.TextColumn(
                                            "Metric",
                                            help="Metric name (read-only)"
                                        ),
                                        "Min": st.column_config.NumberColumn(
                                            "Min",
                                            help="Minimum acceptable value",
                                            format="%.2f"
                                        ),
                                        "Preferred": st.column_config.NumberColumn(
                                            "Preferred",
                                            help="Target/preferred value",
                                            format="%.2f"
                                        ),
                                        "Max": st.column_config.NumberColumn(
                                            "Max",
                                            help="Maximum acceptable value",
                                            format="%.2f"
                                        ),
                                        "Unit": st.column_config.TextColumn(
                                            "Unit",
                                            help="Unit of measurement (read-only)"
                                        ),
                                        "Source": st.column_config.TextColumn(
                                            "Source",
                                            help="Data source - editable"
                                        )
                                    },
                                    key=f"benchmark_editor_{selected_asset}_{selected_subclass}_{group_name}"
                                )

                                # Save button to persist changes
                                if st.button(f"💾 Save Changes", key=f"save_{group_name}"):
                                    # Update session state with overrides
                                    for index, row in edited_df.iterrows():
                                        metric_name = metric_map.get(row['Metric'])
                                        if metric_name:
                                            unit = row['Unit']
                                            # Convert percentages back to decimals
                                            min_val = row['Min'] / 100 if unit == "%" else row['Min']
                                            pref_val = row['Preferred'] / 100 if unit == "%" else row['Preferred']
                                            max_val = row['Max'] / 100 if unit == "%" else row['Max']

                                            # Store override
                                            st.session_state.benchmark_overrides[selected_asset][selected_subclass][metric_name] = [
                                                min_val, pref_val, max_val, row['Source']
                                            ]

                                    st.success(f"✅ Saved {len(edited_df)} benchmark overrides for {group_name}")
                                    st.rerun()

                        else:
                            # VIEW MODE - Display formatted values
                            group_data = []
                            for metric in available_metrics:
                                if metric in selected_benchmarks:
                                    bench_data = selected_benchmarks[metric]
                                    metric_info = get_metric_info(metric)
                                    unit = metric_info.get("unit", "")

                                    if unit == "%":
                                        min_val = f"{bench_data[0]*100:.1f}%"
                                        pref_val = f"{bench_data[1]*100:.1f}%"
                                        max_val = f"{bench_data[2]*100:.1f}%"
                                    elif unit == "x":
                                        min_val = f"{bench_data[0]:.2f}x"
                                        pref_val = f"{bench_data[1]:.2f}x"
                                        max_val = f"{bench_data[2]:.2f}x"
                                    elif unit == "years":
                                        min_val = f"{bench_data[0]:.1f} yrs"
                                        pref_val = f"{bench_data[1]:.1f} yrs"
                                        max_val = f"{bench_data[2]:.1f} yrs"
                                    elif unit in ["$/sf", "$/unit"]:
                                        min_val = f"${bench_data[0]:,.0f}"
                                        pref_val = f"${bench_data[1]:,.0f}"
                                        max_val = f"${bench_data[2]:,.0f}"
                                    else:
                                        min_val = f"{bench_data[0]:,.1f}"
                                        pref_val = f"{bench_data[1]:,.1f}"
                                        max_val = f"{bench_data[2]:,.1f}"

                                    group_data.append({
                                        "Metric": metric.replace("_", " ").title(),
                                        "Min": min_val,
                                        "Preferred": pref_val,
                                        "Max": max_val,
                                        "Source": bench_data[3] if len(bench_data) > 3 else "Industry Standard"
                                    })

                            if group_data:
                                df_group = pd.DataFrame(group_data)

                                # Check if there are any overrides
                                has_overrides = any(metric in overrides for metric in available_metrics)
                                if has_overrides:
                                    st.info("ℹ️ Using customized benchmark values")

                                # Display with custom styling
                                st.dataframe(
                                    df_group,
                                    use_container_width=True,
                                    hide_index=True,
                                    column_config={
                                        "Metric": st.column_config.TextColumn(
                                            "Metric",
                                            help="Click metric names for detailed descriptions"
                                        ),
                                        "Preferred": st.column_config.TextColumn(
                                            "Preferred",
                                            help="Target range for institutional quality"
                                        ),
                                        "Source": st.column_config.TextColumn(
                                            "Source",
                                            help="Data source for benchmark"
                                        )
                                    }
                                )

            # Add detailed metric descriptions section (separate from groups to avoid nesting)
            st.markdown("---")
            with st.expander("📖 Metric Descriptions", expanded=False):
                st.caption("Detailed explanations for all metrics in the selected property type")
                # Get all metrics from all groups
                all_metrics = []
                for group_name, metric_list in metric_groups.items():
                    all_metrics.extend([m for m in metric_list if m in selected_benchmarks])

                # Display descriptions
                for metric in all_metrics:
                    metric_info = get_metric_info(metric)
                    st.markdown(f"**{metric.replace('_', ' ').title()}**")
                    st.caption(f"*{metric_info.get('description', 'No description available')}*")
                    st.info(f"💡 **Why it matters:** {metric_info.get('why_it_matters', 'Metric importance not documented')}")
                    st.markdown("---")

            # Add benchmark comparison tool
            st.markdown("---")
            st.subheader("🔍 Quick Benchmark Check")

            col1, col2, col3 = st.columns(3)
            with col1:
                check_metric = st.selectbox(
                    "Select Metric",
                    options=[m for m in selected_benchmarks.keys()],
                    format_func=lambda x: x.replace("_", " ").title()
                )

            with col2:
                if check_metric:
                    metric_info = get_metric_info(check_metric)
                    unit = metric_info.get("unit", "")

                    if unit == "%":
                        check_value = st.number_input(
                            f"Your Value (%)",
                            min_value=0.0,
                            max_value=100.0,
                            value=5.0,
                            step=0.1
                        ) / 100  # Convert to decimal
                    else:
                        check_value = st.number_input(
                            f"Your Value ({unit})",
                            min_value=0.0,
                            value=1.0,
                            step=0.1
                        )

            with col3:
                if check_metric and check_value:
                    bench_range = selected_benchmarks[check_metric]
                    status = get_status(check_value, bench_range)

                    if status == "OK":
                        st.success(f"✅ {status}: Within preferred range")
                    elif status == "Borderline":
                        st.warning(f"⚠️ {status}: Acceptable but not ideal")
                    else:
                        st.error(f"❌ {status}: Outside acceptable range")

                    st.caption(f"Target: {bench_range[0]:.2f} - {bench_range[2]:.2f}")
                    st.caption(f"Source: {bench_range[3] if len(bench_range) > 3 else 'Industry'}")

        else:
            st.info("Select an asset class and property type to view relevant benchmarks")

    with tab3:
        st.header("📋 Due Diligence Checklist")

        # Default DD Checklist organized by category
        default_dd_categories = {
            "📊 Financial Due Diligence": [
                "T-12 and T-3 operating statements with GL tie-out",
                "Current rent roll with lease abstracts",
                "Accounts receivable aging report",
                "CAM reconciliation (3 years)",
                "Real estate tax bills and assessment history",
                "Insurance policies and loss runs",
                "Utility bills (12 months)",
                "CapEx history (3-5 years)",
                "Property management agreement",
                "Service contracts inventory"
            ],
            "⚖️ Legal & Title": [
                "ALTA survey with Table A items",
                "Title commitment with all exceptions",
                "Tenant estoppel certificates",
                "SNDAs (Subordination, Non-Disturbance Agreements)",
                "All leases and amendments",
                "Operating agreements/CC&Rs/REAs",
                "Zoning confirmation letter",
                "Certificate of occupancy",
                "Business licenses",
                "Litigation search and disclosure"
            ],
            "🏗️ Physical & Environmental": [
                "Property Condition Assessment (PCA)",
                "Structural engineering report",
                "MEP systems evaluation",
                "Roof inspection and warranty",
                "Elevator inspection certificates",
                "Fire/Life safety inspection",
                "ADA compliance assessment",
                "Phase I Environmental Site Assessment",
                "Mold and indoor air quality report",
                "Asbestos and lead paint surveys"
            ],
            "📈 Market & Competitive": [
                "Market comparable lease analysis",
                "Submarket vacancy and absorption trends",
                "Development pipeline (3-mile radius)",
                "Broker opinion of value (BOV)",
                "Tenant demand and tour activity",
                "Employment and demographic analysis",
                "Trade area analysis (retail)",
                "Competitive set benchmarking"
            ],
            "💰 Debt & Capital Structure": [
                "Existing loan documents review",
                "Covenant compliance certificates",
                "Rate cap confirmation and valuation",
                "Prepayment and defeasance analysis",
                "Payoff letters from current lender",
                "UCC and lien searches",
                "Reserve account statements",
                "JV/LP agreement review (if applicable)"
            ]
        }

        # Initialize custom DD items in session state
        if 'custom_dd_items' not in st.session_state:
            st.session_state.custom_dd_items = {}

        if 'items_to_delete' not in st.session_state:
            st.session_state.items_to_delete = set()

        # Merge default and custom DD items
        dd_categories = {}
        for category, items in default_dd_categories.items():
            # Start with default items
            merged_items = list(items)
            # Add custom items for this category
            if category in st.session_state.custom_dd_items:
                merged_items.extend(st.session_state.custom_dd_items[category])
            dd_categories[category] = merged_items

        # Display DD checklist with expanders
        for category, items in dd_categories.items():
            with st.expander(category, expanded=False):
                for item in items:
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.checkbox(item, key=f"dd_{item}")
                    with col2:
                        st.selectbox("", ["Pending", "In Progress", "Complete", "N/A"],
                                   key=f"status_{item}", label_visibility="collapsed")

        # Progress tracker
        st.markdown("---")
        total_items = sum(len(items) for items in dd_categories.values())
        custom_count = sum(len(items) for items in st.session_state.custom_dd_items.values())
        if custom_count > 0:
            st.info(f"Total DD Items: {total_items} ({custom_count} custom)")
        else:
            st.info(f"Total DD Items: {total_items}")

        # Edit Due Diligence Items Section
        st.markdown("---")
        with st.expander("✏️ Edit Due Diligence Items", expanded=False):
            st.markdown("**Add or remove custom due diligence items for each category**")

            # Reset to Defaults button at top
            col1, col2 = st.columns([3, 1])
            with col2:
                if st.button("🔄 Reset to Defaults", use_container_width=True, type="secondary"):
                    st.session_state.custom_dd_items = {}
                    st.session_state.items_to_delete = set()
                    st.success("✅ Reset to default DD items")
                    st.rerun()

            st.markdown("---")

            # Iterate through each category
            for category in default_dd_categories.keys():
                st.markdown(f"### {category}")

                # Show existing items with delete option
                current_items = dd_categories.get(category, [])
                default_items = default_dd_categories.get(category, [])
                custom_items = st.session_state.custom_dd_items.get(category, [])

                if custom_items:
                    st.markdown("**Custom Items (check to delete):**")
                    items_to_remove = []
                    for custom_item in custom_items:
                        if st.checkbox(f"🗑️ {custom_item}", key=f"delete_{category}_{custom_item}"):
                            items_to_remove.append(custom_item)

                    if items_to_remove and st.button(f"Delete Selected Items", key=f"delete_btn_{category}"):
                        # Remove selected items
                        for item in items_to_remove:
                            if category in st.session_state.custom_dd_items:
                                if item in st.session_state.custom_dd_items[category]:
                                    st.session_state.custom_dd_items[category].remove(item)
                                # Clean up empty categories
                                if not st.session_state.custom_dd_items[category]:
                                    del st.session_state.custom_dd_items[category]
                        st.success(f"✅ Deleted {len(items_to_remove)} item(s)")
                        st.rerun()

                # Add new item
                st.markdown("**Add New Item:**")
                col1, col2 = st.columns([3, 1])
                with col1:
                    new_item = st.text_input(
                        "New Item",
                        key=f"new_item_{category}",
                        placeholder="Enter new due diligence item",
                        label_visibility="collapsed"
                    )
                with col2:
                    if st.button("➕ Add", key=f"add_btn_{category}", use_container_width=True):
                        if new_item and new_item.strip():
                            # Initialize category if not exists
                            if category not in st.session_state.custom_dd_items:
                                st.session_state.custom_dd_items[category] = []
                            # Add item if not duplicate
                            if new_item not in current_items:
                                st.session_state.custom_dd_items[category].append(new_item)
                                st.success(f"✅ Added: {new_item}")
                                st.rerun()
                            else:
                                st.warning("⚠️ Item already exists")
                        else:
                            st.warning("⚠️ Please enter an item")

                st.markdown("---")

    with tab4:
        st.header("📄 Report Generation")

        # Check if we have analysis data
        if not hasattr(st.session_state, 'analysis_data') or not st.session_state.analysis_data:
            st.warning("⚠️ Please complete the analysis in the Analysis tab first to generate reports.")
        else:
            analysis_data = st.session_state.analysis_data

            col1, col2, col3 = st.columns(3)
            with col1:
                if st.button("📑 Generate PDF Report", use_container_width=True):
                    pdf_bytes = generate_pdf_report(analysis_data)
                    st.download_button(
                        label="📥 Download PDF",
                        data=pdf_bytes,
                        file_name=f"DealGenie_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                        mime="application/pdf"
                    )
            with col2:
                if st.button("📊 Export to Excel", use_container_width=True):
                    excel_bytes = generate_excel_export(analysis_data)
                    st.download_button(
                        label="📥 Download Excel",
                        data=excel_bytes,
                        file_name=f"DealGenie_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
            with col3:
                if st.button("📈 Export Charts", use_container_width=True):
                    chart_bytes = generate_chart_export(analysis_data)
                    st.download_button(
                        label="📥 Download Charts",
                        data=chart_bytes,
                        file_name=f"DealGenie_Charts_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
                        mime="image/png"
                    )

    with tab5:
        st.header("⚙️ Template Settings & Management")

        # Get list of templates
        available_templates = list_templates()
        default_template = get_default_template()

        # Section 1: Template List
        st.markdown("### 📋 Saved Templates")

        if available_templates:
            # Build dataframe with template info
            template_rows = []
            for template_name in available_templates:
                metadata = get_template_metadata(template_name)
                if metadata:
                    created = metadata.get("created_date", "Unknown")
                    if created != "Unknown":
                        try:
                            created = datetime.fromisoformat(created).strftime("%Y-%m-%d %H:%M")
                        except:
                            pass

                    is_default = "⭐ Yes" if template_name == default_template else "No"

                    template_rows.append({
                        "Template Name": template_name,
                        "Created": created,
                        "Benchmarks": metadata.get("num_benchmarks", 0),
                        "DD Items": metadata.get("num_dd_items", 0),
                        "Default": is_default
                    })

            # Display as dataframe
            df_templates = pd.DataFrame(template_rows)
            st.dataframe(df_templates, use_container_width=True, hide_index=True)

            st.markdown("---")

            # Section 2: Template Actions
            st.markdown("### 🔧 Template Actions")

            col1, col2 = st.columns(2)

            with col1:
                selected_for_action = st.selectbox(
                    "Select Template",
                    options=available_templates,
                    key="template_action_selector"
                )

            with col2:
                st.markdown("**Quick Actions:**")

            # Action buttons
            col1, col2, col3, col4 = st.columns(4)

            with col1:
                if st.button("📂 Load", use_container_width=True, key="settings_load"):
                    if load_template(selected_for_action):
                        st.success(f"✅ Loaded '{selected_for_action}'")
                        st.rerun()

            with col2:
                if st.button("🗑️ Delete", use_container_width=True, key="settings_delete", type="secondary"):
                    if delete_template(selected_for_action):
                        st.success(f"✅ Deleted '{selected_for_action}'")
                        # Clear default if it was the default template
                        if selected_for_action == default_template:
                            clear_default_template()
                        st.rerun()

            with col3:
                # Duplicate button with input
                if st.button("📋 Duplicate", use_container_width=True, key="settings_duplicate"):
                    st.session_state.show_duplicate_input = True

            with col4:
                # Set as default button
                if selected_for_action == default_template:
                    if st.button("⭐ Clear Default", use_container_width=True, key="settings_clear_default"):
                        if clear_default_template():
                            st.success("✅ Cleared default template")
                            st.rerun()
                else:
                    if st.button("⭐ Set as Default", use_container_width=True, key="settings_set_default"):
                        if set_default_template(selected_for_action):
                            st.success(f"✅ Set '{selected_for_action}' as default")
                            st.rerun()

            # Duplicate input (shown when duplicate button clicked)
            if st.session_state.get('show_duplicate_input', False):
                st.markdown("---")
                col1, col2 = st.columns([3, 1])
                with col1:
                    duplicate_name = st.text_input(
                        "New Template Name",
                        placeholder=f"{selected_for_action}_copy",
                        key="duplicate_name_input"
                    )
                with col2:
                    st.markdown("")
                    st.markdown("")
                    if st.button("✅ Create Duplicate", use_container_width=True):
                        if duplicate_name and duplicate_name.strip():
                            if duplicate_template(selected_for_action, duplicate_name):
                                st.success(f"✅ Created '{duplicate_name}'")
                                st.session_state.show_duplicate_input = False
                                st.rerun()
                        else:
                            st.warning("Please enter a name")

        else:
            st.info("No templates saved yet. Create templates using the sidebar or by saving your current settings.")

        # Section 3: Create New Template
        st.markdown("---")
        st.markdown("### ➕ Create New Template")

        with st.expander("Create Template from Scratch", expanded=False):
            col1, col2 = st.columns(2)

            with col1:
                new_template_name = st.text_input(
                    "Template Name",
                    placeholder="e.g., Office Class A CBD",
                    key="new_template_name"
                )

                new_template_description = st.text_area(
                    "Description (optional)",
                    placeholder="Describe the purpose of this template",
                    key="new_template_desc"
                )

            with col2:
                new_template_asset_class = st.selectbox(
                    "Base Asset Class",
                    options=["", "Multifamily", "Office", "Industrial", "Retail", "Hotel"],
                    key="new_template_asset"
                )

                new_template_subclass = st.selectbox(
                    "Property Type",
                    options=[""] + (list(BENCHMARK_DATA.get(new_template_asset_class.lower(), {}).keys()) if new_template_asset_class else []),
                    format_func=lambda x: x.replace("_", " ").title() if x else "Select asset class first",
                    key="new_template_subclass"
                )

            if st.button("➕ Create Template", use_container_width=True, key="create_new_template"):
                if new_template_name and new_template_name.strip():
                    # Create template with specified metadata
                    template_data = {
                        "template_name": new_template_name,
                        "created_date": datetime.now().isoformat(),
                        "description": new_template_description,
                        "benchmark_overrides": {},
                        "custom_dd_items": {},
                        "profile_name": st.session_state.get("profile_name", ""),
                        "asset_class_defaults": {
                            "asset_class": new_template_asset_class,
                            "subclass": new_template_subclass
                        }
                    }

                    if save_template(new_template_name, template_data):
                        st.success(f"✅ Created template '{new_template_name}'")
                        st.rerun()
                else:
                    st.warning("Please enter a template name")

        # Section 4: Template Comparison Tool
        st.markdown("---")
        st.markdown("### 🔍 Template Comparison")

        if len(available_templates) >= 2:
            with st.expander("Compare Two Templates", expanded=False):
                col1, col2 = st.columns(2)

                with col1:
                    compare_template1 = st.selectbox(
                        "First Template",
                        options=available_templates,
                        key="compare_template1"
                    )

                with col2:
                    compare_template2 = st.selectbox(
                        "Second Template",
                        options=[t for t in available_templates if t != compare_template1],
                        key="compare_template2"
                    )

                if st.button("🔍 Compare Templates", use_container_width=True):
                    comparison = compare_templates(compare_template1, compare_template2)

                    if "error" in comparison:
                        st.error(comparison["error"])
                    else:
                        st.markdown(f"**Comparing:** `{comparison['template1']}` vs `{comparison['template2']}`")

                        col1, col2 = st.columns(2)

                        with col1:
                            st.markdown("**📊 Benchmark Differences:**")
                            for diff in comparison['benchmark_diffs']:
                                st.caption(diff)

                        with col2:
                            st.markdown("**📋 DD Item Differences:**")
                            for diff in comparison['dd_diffs']:
                                st.caption(diff)
        else:
            st.info("Create at least 2 templates to use the comparison tool")

        # Section 5: Default Template Info
        st.markdown("---")
        st.markdown("### ⭐ Default Template")

        if default_template:
            st.info(f"**Current Default:** {default_template}")
            st.caption("This template will automatically load when the app starts")
        else:
            st.info("No default template set. You can set one by selecting a template above and clicking 'Set as Default'")

if __name__ == "__main__":
    main()